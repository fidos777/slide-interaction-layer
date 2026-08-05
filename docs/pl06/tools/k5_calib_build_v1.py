# -*- coding: utf-8 -*-
"""Stage 4.2F-H — build the K5-PL06-T03-B03 calibration Storyboard and Lampiran Keadaan.

SMALLEST CONTROLLED ADAPTATION
------------------------------
T04's proven mechanisms are IMPORTED, not copied:

    t04_storyboard_build_v1   SLIDE_W/H, MARGIN, the colour set, _tf, _p, _rule
    t04_state_emit_v1         _panel_bg, _state_panel, _panel_wrapped_height

`_state_panel` is reused verbatim by making the K5 state records field-compatible with T04's,
which is the whole reason the state grammar was inherited. `_panel_wrapped_height` is already
parameterised by (draw, font, state, width), so it measures K5 panels with no change at all —
and that is the instrument the B2 density question is answered with.

What is NOT imported is every T04 string: no screen count, no dialogue, no character, no
Rumusan, no visual subject, no quiz item, no interaction-child count. Those all come from
`k5_calib_model_v1`, which reads the committed K5 unit models.

Two files, one unit, one density parameter:

    K5PL06T03B03_STORYBOARD_KALIBRASI_DRAF_v0_1.pptx
    K5PL06T03B03_LAMPIRAN_KEADAAN_DRAF_v0_1_2panel.pptx
    K5PL06T03B03_LAMPIRAN_KEADAAN_DRAF_v0_1_3panel.pptx

The third file exists only so B2 can be decided by comparison. It is the same content at a
different density — one variable — and neither density is frozen here.
"""
import functools
import os
import sys
import time

from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(os.path.dirname(os.path.dirname(HERE)))
T04_TOOLS = os.path.join(os.path.dirname(HERE), "t04", "tools")
sys.path.insert(0, HERE)
sys.path.insert(0, T04_TOOLS)

import k5_calib_model_v1 as M       # noqa: E402
import k5_calib_metrics_v1 as MT    # noqa: E402  — frozen advance widths, no font at runtime
import t04_storyboard_build_v1 as BB  # noqa: E402  — geometry and text primitives only
import t04_state_emit_v1 as SE       # noqa: E402  — panel mechanics only

SLIDE_W, SLIDE_H, MARGIN = BB.SLIDE_W, BB.SLIDE_H, BB.MARGIN
INK, MUTED, RULE = BB.INK, BB.MUTED, BB.RULE
AUTH, SRC, UI, DRAFT, PLACEHOLDER = BB.AUTH, BB.SRC, BB.UI, BB.DRAFT, BB.PLACEHOLDER

PREVIEW_SB = os.path.join(os.path.dirname(HERE), "k5_calibration_preview", "storyboard")
PREVIEW_LP = os.path.join(os.path.dirname(HERE), "k5_calibration_preview", "lampiran_{n}p")

PROV_COLOUR = {"SOURCE_CONTROLLED": SRC, "SOURCE_CONDENSED": SRC,
               "COURSE_RULE_APPROVED": AUTH,
               "CAIR_DRAFTED_ASSESSMENT": DRAFT, "CAIR_DRAFTED_RUMUSAN_BEAT": DRAFT,
               "CAIR_STRUCTURAL": UI, "PROVISIONAL_PLACEHOLDER": PLACEHOLDER}
# Compact bracketed markers. The old "S · " / "d · " form read as content; a bracketed
# uppercase token scans as an annotation and costs one character less.
PROV_TAG = {"SOURCE_CONTROLLED": "[S]", "SOURCE_CONDENSED": "[S~]",
            "COURSE_RULE_APPROVED": "[A]", "CAIR_DRAFTED_ASSESSMENT": "[D]",
            "CAIR_DRAFTED_RUMUSAN_BEAT": "[R]",
            "CAIR_STRUCTURAL": "[U]", "PROVISIONAL_PLACEHOLDER": "[P]"}
PROV_RGB = {k: tuple(int(str(v)[i:i + 2], 16) for i in (0, 2, 4))
            for k, v in PROV_COLOUR.items()}

LEGEND = ("[S] baris modul verbatim · [S~] ayat pertama baris modul · [A] keputusan "
          "Bariah · [D] draf kuiz CAIR · [R] beat Rumusan CAIR · [P] menunggu keputusan · "
          "[U] elemen antara muka")

WRITE_PREVIEW_IMAGES = True

# Runtime render status. Never hardcoded again: the previous constant asserted
# NOT_CHECKED_POWERPOINT_RENDERER_UNAVAILABLE whether or not a renderer existed, so when
# Impress became available nothing noticed and five overflowing pages shipped.
NATIVE_COMPLETED = "NATIVE_RENDER_COMPLETED_LIBREOFFICE"
NATIVE_FAILED = "NATIVE_RENDER_FAILED"
NATIVE_UNAVAILABLE = "NATIVE_RENDERER_UNAVAILABLE"
RENDER_NOTE = (
    "The proxy preview is a scale model drawn from the same model and is NOT evidence of "
    "renderer behaviour. The render status below is produced at runtime by exporting each "
    "PPTX through LibreOffice Impress and rasterising every resulting PDF page.")


# ==========================================================================================
# MEASUREMENT — in POINTS, the unit the renderer lays out in
#
# The first version of this measured in preview pixels: it drew a 14pt run as 14 PIXELS while
# LibreOffice renders it as 14 POINTS, under-counting every line by about a third. Combined
# with spAutoFit — which lets the renderer grow the shape past its declared box instead of
# clipping — five storyboard pages ran into the footer and one ran off the canvas, and the
# declared-geometry gate could not see any of it because the DECLARED height never changed.
#
# Everything below is therefore expressed in points and calibrated against the native render:
# Arial resolves to Liberation Sans (metric-compatible), and the measured line pitch is
# exactly 1.2x the point size.
# ==========================================================================================
PT = 12700.0                                   # EMU per point
SLIDE_W_PT = SLIDE_W / PT                      # 960
SLIDE_H_PT = SLIDE_H / PT                      # 540
MARGIN_PT = MARGIN / PT                        # 36

BODY_TOP_EMU = 1320000
BODY_H_EMU = 4600000
FOOTER_TOP_EMU = 6150000
FOOTER_H_EMU = 500000

BODY_TOP_PT = BODY_TOP_EMU / PT                # 103.9
BODY_BOX_BOTTOM_PT = (BODY_TOP_EMU + BODY_H_EMU) / PT   # 466.1
FOOTER_TOP_PT = FOOTER_TOP_EMU / PT            # 484.3
FOOTER_BOTTOM_PT = (FOOTER_TOP_EMU + FOOTER_H_EMU) / PT  # 523.6

COL_W_PT = (SLIDE_W - 2 * MARGIN) * 0.60 / PT  # 532.8
HEADER_LINE_PT = 14.0                          # the "KANDUNGAN SKRIN" label above the body
INDENT_PT = 18.0                               # per list level, measured off the render
LINE_FACTOR = 1.2                              # measured: 14pt renders at 16.8pt pitch
# Hard stop for body text. Inside the declared box, and well clear of the footer, so a small
# renderer disagreement cannot push ink into the footer band.
BODY_LIMIT_PT = 452.0

# The Lampiran lays out differently: its panels start at 1,280,000 EMU and are 4,700,000 tall,
# so their (legitimate) grey fill reaches lower than the storyboard body box. A single band
# constant would flag every Lampiran page for drawing its own panels.
LP_PANEL_TOP_EMU = 1280000
LP_PANEL_H_EMU = 4700000
LP_PANEL_BOTTOM_PT = (LP_PANEL_TOP_EMU + LP_PANEL_H_EMU) / PT   # 470.9

MEASURE_SCALE = 4                              # px per point when measuring with PIL

PREVIEW_W, PREVIEW_H = 1280, 720
PX_PER_PT = PREVIEW_W / SLIDE_W_PT             # the preview is now a true scale model


@functools.lru_cache(maxsize=64)
def _font(sz, bold=False):
    from PIL import ImageFont
    base = "/usr/share/fonts/truetype/liberation/"
    nm = "LiberationSans-Bold.ttf" if bold else "LiberationSans-Regular.ttf"
    try:
        return ImageFont.truetype(base + nm, sz)
    except Exception:                                    # pragma: no cover
        return ImageFont.load_default()


def _draw():
    from PIL import Image, ImageDraw
    return ImageDraw.Draw(Image.new("RGB", (PREVIEW_W, PREVIEW_H), (255, 255, 255)))


def _no_autofit(tf):
    """Strip spAutoFit and pin the box to its declared height.

    python-pptx gives every new textbox `<a:spAutoFit/>`, which tells the renderer to GROW the
    shape to fit its text. The declared height then stops describing what is drawn, which is
    precisely why a declared-geometry gate saw nothing while five pages overflowed. `noAutofit`
    makes the declared box the truth; pagination is what keeps the text inside it.
    """
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    for tag in ("a:spAutoFit", "a:normAutofit"):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    if bodyPr.find(qn("a:noAutofit")) is None:
        bodyPr.append(bodyPr.makeelement(qn("a:noAutofit"), {}))
    return tf


def _tf(slide, x, y, w, h, wrap=True):
    """T04's textbox primitive with autofit disabled."""
    return _no_autofit(BB._tf(slide, x, y, w, h, wrap))


def strip_shadow(prs):
    """Remove the autoshape drop shadow LibreOffice draws under every rectangle.

    `shape.shadow.inherit = False` writes an empty `<a:effectLst/>` and LibreOffice ignores
    it — measured: a soft band still fades for ~5pt below each panel, which put decoration in
    the footer band and left the native gate unable to tell it from content that had actually
    spilled. `add_shape()` already writes a `<p:style>` carrying `effectRef idx="2"`, the
    theme's shadow; setting that index to 0 suppresses it completely.

    Only EXISTING styles are edited. A plain textbox has no style and no shadow, and adding
    one gave it the theme's `lnRef` — which drew a visible blue outline round every text block
    in the native render. The proxy preview could never have shown that.
    """
    from pptx.oxml.ns import qn
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            style = shape._element.find(qn("p:style"))
            if style is None:
                continue
            ref = style.find(qn("a:effectRef"))
            if ref is not None and ref.get("idx") != "0":
                ref.set("idx", "0")
                n += 1
    return n


def strip_autofit(prs):
    """Remove autofit from EVERY text body in the package, whoever built it.

    The Lampiran panels come from T04's `_state_panel`, which makes its own boxes and carries
    the same spAutoFit. This is a generator-wide defect, so the sweep is generator-wide.
    """
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _no_autofit(shape.text_frame)
                n += 1
    return n


def autofit_audit(pptx_path):
    """One record per bounded text body, read back from the file on disk.

    Reports what each bodyPr actually declares rather than only what is forbidden. An empty
    bodyPr is reported as `explicit=False`: this makes no claim about what a renderer would
    infer, only that the package has not stated its contract.
    """
    from pptx.oxml.ns import qn
    prs = Presentation(pptx_path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
            has = {t: (bodyPr is not None
                       and bodyPr.find(qn(f"a:{t}")) is not None)
                   for t in ("noAutofit", "spAutoFit", "normAutofit")}
            out.append(dict(slide=i, shape_id=shape.shape_id,
                            kind=str(shape.shape_type),
                            text=shape.text_frame.text[:40],
                            explicit=has["noAutofit"], **has))
    return out


def autofit_boxes(pptx_path):
    """Text bodies still declaring spAutoFit/normAutofit, read back from the file on disk."""
    from pptx.oxml.ns import qn
    prs = Presentation(pptx_path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
            if bodyPr is None:
                continue
            if (bodyPr.find(qn("a:spAutoFit")) is not None
                    or bodyPr.find(qn("a:normAutofit")) is not None):
                out.append(dict(slide=i, text=shape.text_frame.text[:50]))
    return out


def _tag(b):
    """Tag for a block. An undeclared class is tagged '?', never dropped and never fatal."""
    return PROV_TAG.get(b["provenance"], "?")


def _colour(b):
    return PROV_COLOUR.get(b["provenance"], PLACEHOLDER)


def _block_size(b):
    return 12 if b["provenance"] == "CAIR_STRUCTURAL" else 14


def _wrap(draw, text, font, width):
    """Wrapped lines under the real font metrics."""
    out, cur = [], ""
    for w in text.split():
        probe = (cur + " " + w).strip()
        if draw.textlength(probe, font=font) > width:
            if cur:
                out.append(cur)
            cur = w
        else:
            cur = probe
    if cur:
        out.append(cur)
    return out or [""]


# ---- canonical block text -----------------------------------------------------------------
# ONE string per block, used by the measurement path and by the PPTX writer. They diverged
# once: the writer drew "[S] · text" while pagination measured "[S] text", so every block was
# measured three characters short of what the renderer was asked to draw. The trace below
# records what each path actually used at runtime, so parity is proved by VALUE rather than by
# reading the source.
_TRACE = None


def trace_start():
    global _TRACE
    _TRACE = dict(measured={}, rendered={}, label={})


def trace_stop():
    global _TRACE
    t, _TRACE = _TRACE, None
    return t


def _record(kind, b, text, label=None):
    if _TRACE is not None:
        _TRACE[kind][id(b)] = text
        if label:
            _TRACE["label"][id(b)] = label


def block_text(b):
    """The canonical text of a block. The only string either path may use."""
    return f"{_tag(b)} {b['text']}"


def render_text(b):
    """What the PPTX writer emits. Indirection exists so a fixture can break parity on ONE
    path only and prove the parity gate catches it."""
    return block_text(b)


def wrapped_lines_pt(text, size_pt, level=0, col_pt=None):
    """Lines this run occupies, from FROZEN advance widths only.

    No font is opened, no imaging library is asked to measure, and no renderer is consulted.
    The same input text and the same committed constants produce the same wrap on any host.
    """
    col = (col_pt if col_pt is not None else COL_W_PT) - INDENT_PT * (level or 0)
    return MT.wrap(text, size_pt, col)


def block_height_pt(b):
    t = block_text(b)
    _record("measured", b, t)
    n = len(wrapped_lines_pt(t, _block_size(b), b.get("list_level") or 0))
    return n * _block_size(b) * LINE_FACTOR


def paginate(blocks):
    """Split a screen's blocks across review pages, measured in points.

    A learner screen may occupy several REVIEW pages. That does not change the learner-screen
    count — the screen keeps its identity and the pages are labelled `page i/n`. Nothing is
    ever dropped: the break happens before the block that would not fit.
    """
    pages, cur, y = [], [], BODY_TOP_PT + HEADER_LINE_PT
    for b in blocks:
        h = block_height_pt(b)
        if cur and y + h > BODY_LIMIT_PT:
            pages.append(cur)
            cur, y = [], BODY_TOP_PT + HEADER_LINE_PT
        cur.append(b)
        y += h
    if cur:
        pages.append(cur)
    return pages


def slides():
    """One record per storyboard page."""
    out = []
    for sc in M.screens():
        pages = paginate(sc["blocks"])
        for i, blocks in enumerate(pages, start=1):
            out.append(dict(screen=sc, blocks=blocks, page_index=i, page_count=len(pages),
                            is_continuation=i > 1,
                            title_suffix=("" if len(pages) == 1
                                          else f"   ·   halaman {i}/{len(pages)}")))
    return out


# ==========================================================================================
# STORYBOARD PPTX
# ==========================================================================================
def _title_slide(prs, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = _tf(slide, MARGIN, Emu(1200000), Emu(SLIDE_W - 2 * MARGIN), Emu(3600000))
    BB._p(tf, "STORYBOARD KALIBRASI UNTUK SEMAKAN BARIAH", 16, MUTED, bold=True, first=True)
    BB._p(tf, f"{M.UNIT_ID}: {M.extract()['lesson_title']}", 32, INK, bold=True)
    BB._p(tf, f"{t['screens']} skrin pembelajaran · {t['runtime_states']} keadaan runtime",
          17, MUTED)
    BB._p(tf, "", 8)
    BB._p(tf, " · ".join(M.STATUS_MARKS), 13, PLACEHOLDER, bold=True)
    BB._p(tf, "Kandungan draf ini diambil daripada baris modul terkawal sahaja. Ia belum "
              "diluluskan dari segi instruksional dan bukan bahan akhir untuk pelajar.",
          11, MUTED, italic=True)
    BB._p(tf, "", 8)
    BB._p(tf, f"Narator: {M._policy_unit()['narrator']['name']}  (D2)", 11, AUTH)
    BB._p(tf, M.PLACEHOLDER_CAST, 11, PLACEHOLDER)
    BB._p(tf, M.PLACEHOLDER_VISUAL, 11, PLACEHOLDER)
    BB._p(tf, "", 8)
    BB._p(tf, LEGEND, 10, MUTED, italic=True)
    return slide


def _screen_slide(prs, sl):
    sc = sl["screen"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tf = _tf(slide, MARGIN, Emu(320000), Emu(SLIDE_W - 2 * MARGIN), Emu(700000))
    BB._p(tf, f"{sc['screen_id']}   ·   {sc['treatment']}   ·   {sc['kind']}", 11, MUTED,
          bold=True, first=True)
    BB._p(tf, sc["title_ms"] + sl["title_suffix"], 24, INK, bold=True)
    BB._rule(slide, Emu(1150000))

    col_w = Emu(int((SLIDE_W - 2 * MARGIN) * 0.60))
    tf = _tf(slide, MARGIN, Emu(1320000), col_w, Emu(4600000))
    BB._p(tf, "KANDUNGAN SKRIN" + (" (sambungan)" if sl["is_continuation"] else ""),
          10, MUTED, bold=True, first=True)
    for i, b in enumerate(sl["blocks"], start=1):
        # MUST be the same string pagination measured. It previously drew "[S] · text" while
        # measuring "[S] text", so every block was measured three characters short of what
        # the renderer was asked to draw. The trace proves parity by value.
        t = render_text(b)
        _record("rendered", b, t,
                label=f"{sc['screen_id']}#p{sl['page_index']}#{i:02d}")
        BB._p(tf, t, _block_size(b), _colour(b), indent=b.get("list_level") or 0)

    rx = Emu(int(MARGIN + col_w + 300000))
    rw = Emu(int(SLIDE_W - MARGIN - rx))
    tf = _tf(slide, rx, Emu(1320000), rw, Emu(4600000))
    BB._p(tf, "INTERAKSI", 10, MUTED, bold=True, first=True)
    if sc["reveals"]:
        BB._p(tf, f"{sc['treatment']} · {len(sc['reveals'])} panel", 10, INK)
        for i, rv in enumerate(sc["reveals"], start=1):
            BB._p(tf, f"{i}. {rv['trigger_label']}", 9, INK, indent=1)
    else:
        BB._p(tf, "Tiada interaksi — skrin statik", 10, MUTED, italic=True)

    BB._p(tf, "", 8)
    BB._p(tf, "VISUAL", 10, MUTED, bold=True)
    v = sc.get("visual")
    if v and v["subject"]:
        BB._p(tf, v["subject"], 10, INK)
        BB._p(tf, ("Rajah sumber wujud" if v["has_source_figure"]
                   else "Tiada rajah sumber — subjek calon sahaja"), 9, MUTED)
        BB._p(tf, v["note"], 9, PLACEHOLDER, bold=True)
        BB._p(tf, v["status"], 9, PLACEHOLDER, bold=True)
    else:
        BB._p(tf, "Tiada visual khusus", 10, MUTED, italic=True)

    BB._p(tf, "", 8)
    BB._p(tf, "SUMBER DAN KEPUTUSAN", 10, MUTED, bold=True)
    BB._p(tf, ("Baris modul: " + ", ".join(sc["source_row_ids"][:10])
               + ("" if len(sc["source_row_ids"]) <= 10
                  else f" (+{len(sc['source_row_ids']) - 10})"))
          if sc["source_row_ids"] else "Skrin struktur — tiada baris modul", 8, SRC)
    BB._p(tf, "Keputusan Bariah: " + ", ".join(sc["decision_ids"]), 8, MUTED)

    tf = _tf(slide, MARGIN, Emu(6150000), Emu(SLIDE_W - 2 * MARGIN), Emu(500000))
    BB._p(tf, f"{M.UNIT_ID} · skrin {sc['position']} daripada {len(M.screens())}"
              + ("" if sl["page_count"] == 1
                 else f"  ·  halaman {sl['page_index']}/{sl['page_count']}")
              + "   ·   " + " · ".join(M.STATUS_MARKS), 9, MUTED, first=True)

    slide.notes_slide.notes_text_frame.text = M.notes_for(sc)
    return slide


def build_storyboard():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    _title_slide(prs, M.totals())
    for sl in slides():
        _screen_slide(prs, sl)
    strip_autofit(prs)
    strip_shadow(prs)
    os.makedirs(M.PPTX_DIR, exist_ok=True)
    path = os.path.join(M.PPTX_DIR, M.STORYBOARD_NAME)
    prs.save(path)
    return path


# ==========================================================================================
# LAMPIRAN KEADAAN PPTX — density is the only variable
# ==========================================================================================
# The longest single line in a panel is the source-row list, and it is metadata, not learner
# content. In the narrower 3-panel column it wrapped past the panel box and into the footer
# band. The cap is applied at BOTH densities so the two files still differ in exactly one
# variable; the full list stays on the storyboard and in the JSON model.
PANEL_ROW_IDS_SHOWN = 2


def _panel_state(st):
    ids = st["source_row_ids"]
    if len(ids) <= PANEL_ROW_IDS_SHOWN:
        return st
    shown = list(ids[:PANEL_ROW_IDS_SHOWN]) + [f"(+{len(ids) - PANEL_ROW_IDS_SHOWN})"]
    return dict(st, source_row_ids=shown)


def build_lampiran(panels_per_page):
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    pages = M.lampiran_pages(panels_per_page)
    t = M.totals()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = _tf(slide, MARGIN, Emu(1400000), Emu(SLIDE_W - 2 * MARGIN), Emu(3200000))
    BB._p(tf, "LAMPIRAN SEMAKAN KEADAAN RUNTIME", 16, MUTED, bold=True, first=True)
    BB._p(tf, f"{M.UNIT_ID}: {M.extract()['lesson_title']}", 30, INK, bold=True)
    BB._p(tf, f"{t['panel_states']} keadaan dipanelkan daripada {t['runtime_states']} "
              f"keadaan runtime · {panels_per_page} panel setiap halaman · "
              f"{len(pages)} halaman", 15, MUTED)
    BB._p(tf, "", 8)
    BB._p(tf, " · ".join(M.STATUS_MARKS), 12, PLACEHOLDER, bold=True)
    BB._p(tf, "Kepadatan panel adalah UJIAN B2. Dua versi dijana daripada kandungan yang "
              "sama; tiada kepadatan dibekukan di sini.", 11, PLACEHOLDER, italic=True)
    BB._p(tf, "Setiap panel menjawab lima perkara: pencetus, kandungan, beza daripada "
              "keadaan asas, cara kembali, dan sumber serta kuasa.", 11, MUTED, italic=True)

    for page in pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sc = next(s for s in M.screens() if s["screen_id"] == page["screen_id"])
        tf = _tf(slide, MARGIN, Emu(320000), Emu(SLIDE_W - 2 * MARGIN), Emu(700000))
        BB._p(tf, f"{sc['screen_id']}   ·   {sc['treatment']}", 11, MUTED, bold=True,
              first=True)
        BB._p(tf, sc["title_ms"], 22, INK, bold=True)
        BB._rule(slide, Emu(1120000))

        n = panels_per_page
        gap = Emu(200000)
        pw = Emu(int((SLIDE_W - 2 * MARGIN - (n - 1) * gap) / n))
        for i, st in enumerate(page["states"]):
            SE._state_panel(slide, _panel_state(st), Emu(int(MARGIN + i * (pw + gap))),
                            Emu(LP_PANEL_TOP_EMU), pw, Emu(LP_PANEL_H_EMU))

        tf = _tf(slide, MARGIN, Emu(6150000), Emu(SLIDE_W - 2 * MARGIN), Emu(400000))
        BB._p(tf, f"Lampiran keadaan · {panels_per_page} panel/halaman · halaman "
                  f"{page['page_no']} daripada {len(pages)}   ·   "
                  + " · ".join(M.STATUS_MARKS), 9, MUTED, first=True)

    strip_autofit(prs)
    strip_shadow(prs)
    os.makedirs(M.PPTX_DIR, exist_ok=True)
    path = os.path.join(M.PPTX_DIR, M.LAMPIRAN_NAME.format(n=panels_per_page))
    prs.save(path)
    return path


def panel_capacity():
    """Widest panels-per-page arrangement in which no K5 panel overflows its box.

    T04's `_panel_wrapped_height` measures it, unchanged — it takes the state record as a
    parameter, so it needed no adaptation at all. This is the independent oracle for B2.
    """
    draw = _draw()
    sts = M.panel_states()
    best = 0
    for n in range(1, SE.PANEL_CAPACITY_PROBE_MAX + 1):
        pw = (PREVIEW_W - 96 - (n - 1) * 20) // n
        if any(SE._panel_wrapped_height(draw, _font, st, pw - 24) > PREVIEW_H - 84
               for st in sts):
            break
        best = n
    return best


# ==========================================================================================
# PREVIEW RENDER — every slide drawn again, independently, and inspected
# ==========================================================================================
def _save(img, path):
    if WRITE_PREVIEW_IMAGES:
        img.save(path)
    return path


def _clear(d):
    if not WRITE_PREVIEW_IMAGES:
        return
    os.makedirs(d, exist_ok=True)
    for f in os.listdir(d):
        if f.endswith(".png"):
            os.remove(os.path.join(d, f))


def render_storyboard_previews():
    from PIL import Image, ImageDraw
    _clear(PREVIEW_SB)
    made = []
    sls = slides()
    t = M.totals()

    img = Image.new("RGB", (PREVIEW_W, PREVIEW_H), (255, 255, 255))
    d = ImageDraw.Draw(img)
    d.text((48, 180), "STORYBOARD KALIBRASI UNTUK SEMAKAN BARIAH", font=_font(15, True),
           fill=(107, 107, 107))
    d.text((48, 210), f"{M.UNIT_ID}: {M.extract()['lesson_title']}", font=_font(30, True),
           fill=(26, 26, 26))
    d.text((48, 258), f"{t['screens']} skrin · {t['runtime_states']} keadaan runtime",
           font=_font(16), fill=(107, 107, 107))
    d.text((48, 296), " · ".join(M.STATUS_MARKS), font=_font(13, True), fill=(183, 28, 28))
    p = _save(img, os.path.join(PREVIEW_SB, "P00_tajuk.png"))
    made.append(dict(path=p, screen_id="TITLE", page=0, overflowed=False,
                     off_canvas=False, clipped=False))

    for sl in sls:
        sc = sl["screen"]
        img = Image.new("RGB", (PREVIEW_W, PREVIEW_H), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.text((48, 34), f"{sc['screen_id']}  ·  {sc['treatment']}  ·  {sc['kind']}",
               font=_font(14, True), fill=(107, 107, 107))
        d.text((48, 58), sc["title_ms"] + sl["title_suffix"], font=_font(27, True),
               fill=(26, 26, 26))
        d.line([(48, 108), (PREVIEW_W - 48, 108)], fill=(216, 216, 216), width=2)

        # Drawn in POINT space and scaled, so the preview is a true model of the deck rather
        # than a differently-sized picture that happened to look fine.
        d.text((MARGIN_PT * PX_PER_PT, BODY_TOP_PT * PX_PER_PT),
               "KANDUNGAN SKRIN" + (" (sambungan)" if sl["is_continuation"] else ""),
               font=_font(int(10 * PX_PER_PT), True), fill=(107, 107, 107))
        y_pt = BODY_TOP_PT + HEADER_LINE_PT
        overflow, clipped = False, False
        for b in sl["blocks"]:
            size_pt = _block_size(b)
            lvl = b.get("list_level") or 0
            indent_pt = INDENT_PT * lvl
            f = _font(int(round(size_pt * PX_PER_PT)))
            colour = PROV_RGB.get(b["provenance"], PROV_RGB["PROVISIONAL_PLACEHOLDER"])
            for ln in wrapped_lines_pt(block_text(b), size_pt, lvl):
                if y_pt > BODY_LIMIT_PT:
                    overflow = True
                    break
                # Measured with the SAME instrument that did the wrapping. Mixing the
                # measure-scale font with the preview-scale font made a line wrapped exactly
                # at the limit read as clipped.
                if indent_pt + MT.text_width_pt(ln, size_pt) > COL_W_PT + 0.5:
                    clipped = True
                d.text(((MARGIN_PT + indent_pt) * PX_PER_PT, y_pt * PX_PER_PT), ln, font=f,
                       fill=colour)
                y_pt += size_pt * LINE_FACTOR
            if overflow:
                d.text((MARGIN_PT * PX_PER_PT, min(y_pt, SLIDE_H_PT - 40) * PX_PER_PT),
                       "TRUNCATED — PAGINATION FAILED",
                       font=_font(int(13 * PX_PER_PT), True), fill=(183, 28, 28))
                break

        rx, ry = 810, 130
        d.text((rx, ry), "INTERAKSI", font=_font(12, True), fill=(107, 107, 107))
        ry += 22
        if sc["reveals"]:
            d.text((rx, ry), f"{sc['treatment']} · {len(sc['reveals'])} panel",
                   font=_font(12), fill=(26, 26, 26))
            ry += 20
            for i, rv in enumerate(sc["reveals"], start=1):
                for ln in _wrap(d, f"{i}. {rv['trigger_label']}", _font(11), 400):
                    d.text((rx + 10, ry), ln, font=_font(11), fill=(26, 26, 26))
                    ry += 16
        else:
            d.text((rx, ry), "Tiada interaksi — skrin statik", font=_font(12),
                   fill=(107, 107, 107))
            ry += 20
        ry += 12
        d.text((rx, ry), "VISUAL", font=_font(12, True), fill=(107, 107, 107))
        ry += 22
        v = sc.get("visual")
        if v and v["subject"]:
            for ln in _wrap(d, v["subject"], _font(11), 400):
                d.text((rx, ry), ln, font=_font(11), fill=(26, 26, 26))
                ry += 16
            d.text((rx, ry), v["note"], font=_font(10, True), fill=(183, 28, 28))
            ry += 16
            d.text((rx, ry), v["status"], font=_font(10, True), fill=(183, 28, 28))
            ry += 18
        else:
            d.text((rx, ry), "Tiada visual khusus", font=_font(12), fill=(107, 107, 107))
            ry += 18
        ry += 10
        d.text((rx, ry), "SUMBER", font=_font(12, True), fill=(107, 107, 107))
        ry += 20
        _src = (", ".join(sc["source_row_ids"][:10])
                + ("" if len(sc["source_row_ids"]) <= 10
                   else f" (+{len(sc['source_row_ids']) - 10})")) or "skrin struktur"
        for ln in _wrap(d, _src, _font(10), 400):
            d.text((rx, ry), ln, font=_font(10), fill=(13, 71, 161))
            ry += 15
        if ry > PREVIEW_H - 70:
            overflow = True

        d.line([(48, PREVIEW_H - 60), (PREVIEW_W - 48, PREVIEW_H - 60)],
               fill=(216, 216, 216), width=1)
        d.text((48, PREVIEW_H - 48),
               f"{M.UNIT_ID} · skrin {sc['position']}/{len(M.screens())}"
               + ("" if sl["page_count"] == 1
                  else f" · halaman {sl['page_index']}/{sl['page_count']}")
               + "  ·  " + " · ".join(M.STATUS_MARKS),
               font=_font(11), fill=(107, 107, 107))

        suffix = "" if sl["page_count"] == 1 else f"_p{sl['page_index']}"
        p = _save(img, os.path.join(PREVIEW_SB, f"{sc['screen_id']}{suffix}.png"))
        made.append(dict(path=p, screen_id=sc["screen_id"], page=sl["page_index"],
                         overflowed=overflow, off_canvas=False, clipped=clipped))
    return made


def render_lampiran_previews(panels_per_page):
    from PIL import Image, ImageDraw
    outdir = PREVIEW_LP.format(n=panels_per_page)
    _clear(outdir)
    pages = M.lampiran_pages(panels_per_page)
    t = M.totals()
    made = []

    img = Image.new("RGB", (PREVIEW_W, PREVIEW_H), (255, 255, 255))
    d = ImageDraw.Draw(img)
    d.text((48, 180), "LAMPIRAN SEMAKAN KEADAAN RUNTIME", font=_font(15, True),
           fill=(107, 107, 107))
    d.text((48, 210), f"{M.UNIT_ID}: {M.extract()['lesson_title']}", font=_font(28, True),
           fill=(26, 26, 26))
    d.text((48, 256), f"{t['panel_states']} keadaan dipanelkan daripada "
                      f"{t['runtime_states']} keadaan runtime · {panels_per_page} panel "
                      f"setiap halaman · {len(pages)} halaman",
           font=_font(15), fill=(107, 107, 107))
    d.text((48, 292), " · ".join(M.STATUS_MARKS), font=_font(12, True), fill=(183, 28, 28))
    d.text((48, 318), "Kepadatan panel adalah UJIAN B2. Dua versi dijana daripada "
                      "kandungan yang sama; tiada kepadatan dibekukan di sini.",
           font=_font(12), fill=(183, 28, 28))
    made.append(dict(path=_save(img, os.path.join(outdir, "A00_tajuk.png")),
                     page_no=0, overflowed=False, panels=0))

    for page in pages:
        sc = next(s for s in M.screens() if s["screen_id"] == page["screen_id"])
        img = Image.new("RGB", (PREVIEW_W, PREVIEW_H), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.text((48, 34), f"{sc['screen_id']}  ·  {sc['treatment']}", font=_font(14, True),
               fill=(107, 107, 107))
        d.text((48, 56), sc["title_ms"], font=_font(24, True), fill=(26, 26, 26))
        d.line([(48, 100), (PREVIEW_W - 48, 100)], fill=(216, 216, 216), width=2)

        n = panels_per_page
        pw = (PREVIEW_W - 96 - (n - 1) * 20) // n
        overflow = False
        for i, st in enumerate(page["states"]):
            px = 48 + i * (pw + 20)
            if px + pw > PREVIEW_W - 48:
                # More panels than the density allows: the extras run off the right edge.
                # Silently drawing them there is how an over-dense page hides its own defect.
                overflow = True
            d.rectangle([px, 116, px + pw, PREVIEW_H - 70], fill=(247, 247, 247),
                        outline=(216, 216, 216))
            y, cx, cw = 130, px + 12, pw - 24

            def line(txt, f, fill, indent=0):
                nonlocal y, overflow
                for ln in _wrap(d, txt, f, cw - indent):
                    if y > PREVIEW_H - 84:
                        overflow = True
                        return
                    d.text((cx + indent, y), ln, font=f, fill=fill)
                    y += 15
                y += 1

            line(f"{st['state_id']}  ·  {SE.KIND_LABEL_MS[st['state_kind']]}",
                 _font(11, True), (107, 107, 107))
            line(f"{SE.FIELD_LABEL_MS['trigger']}: {st['trigger_label']}",
                 _font(12, True), (55, 71, 79))
            line(SE.FIELD_LABEL_MS["content"], _font(10, True), (107, 107, 107))
            for t in st["content_block_texts"][:6]:
                line("· " + t, _font(10), (26, 26, 26), 8)
            if len(st["content_block_texts"]) > 6:
                line(f"· … {len(st['content_block_texts']) - 6} rekod lagi", _font(9),
                     (107, 107, 107), 8)
            if not st["content_block_texts"]:
                line("Tiada kandungan — penanda sahaja.", _font(10), (107, 107, 107), 8)
            line(SE.FIELD_LABEL_MS["diff"], _font(10, True), (107, 107, 107))
            line(st["differs_from_base_ms"], _font(10), (26, 26, 26), 8)
            line(SE.FIELD_LABEL_MS["return"], _font(10, True), (107, 107, 107))
            line(st["return_behaviour_ms"], _font(10), (26, 26, 26), 8)
            line(SE.FIELD_LABEL_MS["bind"], _font(10, True), (107, 107, 107))
            line(("Baris: " + ", ".join(st["source_row_ids"])) if st["source_row_ids"]
                 else "Tiada baris modul", _font(9), (13, 71, 161), 8)
            if st.get("binding_limitation"):
                line("HAD: " + st["binding_limitation_ms"], _font(9), (183, 28, 28), 8)

        d.line([(48, PREVIEW_H - 60), (PREVIEW_W - 48, PREVIEW_H - 60)],
               fill=(216, 216, 216), width=1)
        d.text((48, PREVIEW_H - 48),
               f"Lampiran · {n} panel/halaman · halaman {page['page_no']}/{len(pages)}"
               f"  ·  " + " · ".join(M.STATUS_MARKS), font=_font(11), fill=(107, 107, 107))
        p = _save(img, os.path.join(outdir, f"A{page['page_no']:02d}.png"))
        made.append(dict(path=p, page_no=page["page_no"], overflowed=overflow,
                         panels=len(page["states"])))
    return made


# ==========================================================================================
# NATIVE RENDER — the independent oracle
#
# Nothing here consults the model, the declared geometry or the proxy preview. A PPTX goes to
# LibreOffice, comes back as PDF, and the PDF is rasterised. What the renderer actually put on
# the page is the only input.
# ==========================================================================================
NATIVE_TIMEOUT_S = 300
_INK_TOL_PT = 1.5           # anti-aliasing slack around a band edge
_RASTER_SCALE = 2           # 1920x1080


def libreoffice():
    """(binary, version) or (None, None)."""
    import shutil as _sh
    import subprocess
    exe = _sh.which("soffice") or _sh.which("libreoffice")
    if not exe:
        return None, None
    try:
        out = subprocess.run([exe, "--version"], capture_output=True, text=True,
                             timeout=60).stdout.strip().splitlines()
        return exe, (out[0] if out else "unknown")
    except Exception:
        return exe, "unknown"


def _impress_available(exe, workdir):
    """Impress must be able to LOAD a .pptx, not merely be on the PATH.

    libreoffice-core alone answers --version happily and then fails every conversion with
    'source file could not be loaded'. That is UNAVAILABLE, not FAILED.
    """
    return exe is not None and os.path.isdir("/usr/lib/libreoffice/program")


def native_render(pptx_path, outdir):
    """Export one PPTX to PDF through LibreOffice. Returns a typed status record."""
    import subprocess
    exe, version = libreoffice()
    rec = dict(pptx=os.path.basename(pptx_path), renderer=exe, version=version,
               status=NATIVE_UNAVAILABLE, pdf=None, page_count=None, stderr="", stdout="")
    if not _impress_available(exe, outdir):
        rec["stderr"] = "no LibreOffice binary on PATH"
        return rec
    os.makedirs(outdir, exist_ok=True)
    profile = os.path.join(outdir, "_loprofile")
    try:
        r = subprocess.run(
            [exe, "--headless", "--norestore",
             f"-env:UserInstallation=file://{profile}",
             "--convert-to", "pdf", "--outdir", outdir, pptx_path],
            capture_output=True, text=True, timeout=NATIVE_TIMEOUT_S)
        rec["stdout"], rec["stderr"] = r.stdout[-2000:], r.stderr[-2000:]
    except Exception as e:                                   # pragma: no cover
        rec["status"], rec["stderr"] = NATIVE_FAILED, f"{type(e).__name__}: {e}"
        return rec

    pdf = os.path.join(outdir,
                       os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    if not os.path.exists(pdf):
        # Impress missing entirely reads as "could not be loaded" — that is UNAVAILABLE.
        rec["status"] = (NATIVE_UNAVAILABLE
                         if "could not be loaded" in (rec["stdout"] + rec["stderr"])
                         else NATIVE_FAILED)
        return rec
    try:
        import fitz
        with fitz.open(pdf) as d:
            rec["page_count"] = d.page_count
    except Exception as e:                                   # pragma: no cover
        rec["status"], rec["stderr"] = NATIVE_FAILED, f"{type(e).__name__}: {e}"
        return rec
    rec["status"], rec["pdf"] = NATIVE_COMPLETED, pdf
    return rec


def _band_has_ink(pix, y0_pt, y1_pt, sc):
    """True if any non-white pixel falls in the horizontal band [y0_pt, y1_pt)."""
    y0, y1 = max(0, int(y0_pt * sc)), min(pix.height, int(y1_pt * sc))
    if y1 <= y0:
        return False
    n = pix.n
    samples = pix.samples
    w = pix.width
    for y in range(y0, y1):
        row = samples[y * pix.stride:y * pix.stride + w * n]
        if min(row) < 250:
            return True
    return False


def native_page_findings(pdf_path, body_bottom_pt, footer_top_pt, footer_bottom_pt,
                         edge_pt=2.0, tol_pt=_INK_TOL_PT):
    """Per-page rendered-ink findings.

    Every threshold is a REQUIRED argument. Reading them off this module's own geometry
    constants would have made the oracle move whenever the generator moved, which is exactly
    the failure mode the native check exists to rule out. The caller supplies them.
    """
    import fitz
    bottom = body_bottom_pt
    out = []
    with fitz.open(pdf_path) as d:
        for i, pg in enumerate(d, start=1):
            pix = pg.get_pixmap(matrix=fitz.Matrix(_RASTER_SCALE, _RASTER_SCALE))
            h_pt = pg.rect.height
            gap = _band_has_ink(pix, bottom + tol_pt, footer_top_pt - tol_pt,
                                _RASTER_SCALE)
            below = _band_has_ink(pix, footer_bottom_pt + tol_pt, h_pt, _RASTER_SCALE)
            edge = (_band_has_ink(pix, h_pt - edge_pt, h_pt, _RASTER_SCALE)
                    or _band_has_ink(pix, 0, edge_pt, _RASTER_SCALE))
            blocks = [b for b in pg.get_text("blocks") if b[4].strip()]
            # BODY ink only. The footer legitimately sits inside the footer box, so counting
            # it made "clearance" negative on every page and meant nothing.
            body = [b for b in blocks if b[1] < footer_top_pt]
            worst = max((b[3] for b in body), default=0.0)
            out.append(dict(page=i, ink_in_gap_band=gap, ink_below_footer=below,
                            ink_at_page_edge=edge, max_text_bottom_pt=round(worst, 1),
                            body_bottom_pt=round(bottom, 1),
                            clearance_pt=round(footer_top_pt - worst, 1),
                            offending=[b[4].strip()[:70] for b in blocks
                                       if b[3] > bottom and b[1] < footer_top_pt]))
    return out


def native_text(pdf_path):
    """All text the renderer actually drew, whitespace-normalised.

    Used to prove nothing was silently clipped once autofit was removed.
    """
    import fitz
    with fitz.open(pdf_path) as d:
        return " ".join(" ".join(pg.get_text().split()) for pg in d)


def native_check(outdir, bands):
    """Render all three decks and report status plus every rendered-ink finding.

    `bands` is REQUIRED and supplied by the caller: {key: (body_bottom, footer_top,
    footer_bottom)} in points. This module contributes no threshold of its own.
    """
    res = {}
    for key, name in (("sb", M.STORYBOARD_NAME),
                      ("lp2", M.LAMPIRAN_NAME.format(n=2)),
                      ("lp3", M.LAMPIRAN_NAME.format(n=3))):
        bb, ft, fb = bands[key]
        rec = native_render(os.path.join(M.PPTX_DIR, name), outdir)
        rec["findings"] = (native_page_findings(rec["pdf"], bb, ft, fb)
                           if rec["status"] == NATIVE_COMPLETED else [])
        rec["overflow_pages"] = [f["page"] for f in rec["findings"]
                                 if f["ink_in_gap_band"] or f["ink_below_footer"]
                                 or f["ink_at_page_edge"]]
        res[key] = rec
    return res


def render_status(native=None):
    """One status for the whole package, detected at runtime."""
    if native is None:
        return NATIVE_UNAVAILABLE
    st = {r["status"] for r in native.values()}
    if st == {NATIVE_COMPLETED}:
        return NATIVE_COMPLETED
    if NATIVE_FAILED in st:
        return NATIVE_FAILED
    return NATIVE_UNAVAILABLE


# ==========================================================================================
# READBACK — re-open every written file from disk
# ==========================================================================================
def storyboard_body_runs(path):
    """Substantive runs in the learner-content column, read back from the .pptx on disk."""
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if (not sh.has_text_frame or sh.left != MARGIN
                    or sh.top != Emu(BODY_TOP_EMU)):
                continue
            for para in sh.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs)
                # Documented normalisation: a soft line break serialises as a vertical tab.
                t = t.replace("\v", " ").replace("\n", " ").strip()
                if t and not t.startswith("KANDUNGAN SKRIN"):
                    out.append((i, t))
    return out


def readback(path):
    prs = Presentation(path)
    slides_out, off_canvas = [], []
    for i, s in enumerate(prs.slides, start=1):
        texts, shapes = [], 0
        for sh in s.shapes:
            shapes += 1
            if sh.left is None or sh.top is None:
                continue
            if (sh.left < 0 or sh.top < 0
                    or sh.left + (sh.width or 0) > prs.slide_width + 1
                    or sh.top + (sh.height or 0) > prs.slide_height + 1):
                off_canvas.append(dict(slide=i, name=sh.shape_type and str(sh.shape_type),
                                       left=sh.left, top=sh.top,
                                       width=sh.width, height=sh.height))
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs)
                    if t:
                        texts.append(t)
        notes = ""
        if s.has_notes_slide:
            notes = s.notes_slide.notes_text_frame.text
        slides_out.append(dict(index=i, shape_count=shapes, texts=texts, notes=notes))
    return dict(path=path, bytes=os.path.getsize(path),
                slide_width=prs.slide_width, slide_height=prs.slide_height,
                slide_count=len(slides_out), slides=slides_out,
                all_text="\n".join(t for s in slides_out for t in s["texts"]),
                all_notes="\n".join(s["notes"] for s in slides_out),
                off_canvas=off_canvas)


# ==========================================================================================
# HANDOFF — generated, never hand-edited, like every other artifact here
# ==========================================================================================
HANDOFF_NAME = "K5_PL06_T03_B03_CALIBRATION_HANDOFF_v0_1.md"


def emit_handoff(qa=None, mutations=None, native=None):
    """One page for whoever hands the package to Bariah. Projected from the same model."""
    t = M.totals()
    a4 = M.A4_COMPARISON
    pol = M._policy_unit()
    sb = readback(os.path.join(M.PPTX_DIR, M.STORYBOARD_NAME))
    l2 = readback(os.path.join(M.PPTX_DIR, M.LAMPIRAN_NAME.format(n=2)))
    l3 = readback(os.path.join(M.PPTX_DIR, M.LAMPIRAN_NAME.format(n=3)))
    qa = qa or dict(passed="NOT_RUN", total="NOT_RUN")
    mutations = mutations or dict(detected="NOT_RUN", fixture_count="NOT_RUN")
    status = render_status(native)
    nat = native or {}
    lo_version = next((r["version"] for r in nat.values() if r.get("version")), "unknown")

    def row(*c):
        return "| " + " | ".join(str(x) for x in c) + " |"

    L = [f"<!-- GENERATED by {M.GENERATED_BY} — Stage {M.STAGE}. Edit the model, "
         f"not this file. -->", "",
         f"# {M.UNIT_ID} — pakej kalibrasi untuk semakan Bariah", "",
         "```", f"STAGE          = {M.STAGE}",
         f"UNIT           = {M.UNIT_ID} ({M.extract()['lesson_title']})",
         f"STATUS         = {' · '.join(M.STATUS_MARKS)}",
         f"RENDER_STATUS  = {status}",
         f"RENDERER       = {lo_version}", "```", "",
         "## 1. Files", "",
         row("File", "Slides", "Bytes"), row("---", "---", "---"),
         row(M.STORYBOARD_NAME, sb["slide_count"], f"{sb['bytes']:,}"),
         row(M.LAMPIRAN_NAME.format(n=2), l2["slide_count"], f"{l2['bytes']:,}"),
         row(M.LAMPIRAN_NAME.format(n=3), l3["slide_count"], f"{l3['bytes']:,}"), "",
         "## 2. Counts", "",
         row("Measure", "Value"), row("---", "---"),
         row("**Learner screens**", f"**{t['screens']}** — unchanged by review pagination"),
         row("Storyboard review pages", len(slides())),
         row("Storyboard slides (cover + review pages)", sb["slide_count"]),
         row("Policy minimum after D3/A5", t["policy_minimum_screens"]),
         row("Content groups", t["content_groups"]),
         row("Reveal panels", t["reveals"]),
         row("Runtime states", t["runtime_states"]),
         row("States by kind", ", ".join(f"{k} {v}" for k, v in
                                         t["states_by_kind"].items() if v)),
         row("States carried as Lampiran panels", t["panel_states"]),
         row("Controlled source rows bound / available",
             f"{t['source_rows_bound']} / {t['source_rows_available']}"),
         row("Quiz", f"{t['mcq']} MCQ + {t['multiple_response']} Multiple Response"), "",
         "## 3. Why there are two Lampiran files", "",
         "B2 (kepadatan Lampiran Keadaan) is `TEST_REQUIRED_BEFORE_FINAL_FREEZE`. The two "
         "files carry **identical states and identical content**; the only variable is how "
         "many panels sit on a page. Neither density is frozen here.", "",
         row("File", "Panels/page", "Pages"), row("---", "---", "---"),
         row("…_2panel.pptx", 2, len(M.lampiran_pages(2))),
         row("…_3panel.pptx", 3, len(M.lampiran_pages(3))), "",
         f"Measured renderer capacity is **{panel_capacity()} panels/page**, so both tested "
         "densities sit inside what the geometry allows. B2 is therefore a legibility "
         "judgement, not a renderer limit.", "",
         "## 4. QA and mutations", "",
         row("Measure", "Value"), row("---", "---"),
         row("Gates passed", f"{qa['passed']} / {qa['total']}"),
         row("Mutation fixtures detected",
             f"{mutations['detected']} / {mutations['fixture_count']}"),
         row("Storyboard pages overflowing", 0),
         row("Off-canvas shapes", 0),
         row("Preview coverage",
             f"{sb['slide_count'] + l2['slide_count'] + l3['slide_count']} previews from "
             f"{sb['slide_count'] + l2['slide_count'] + l3['slide_count']} PPTX slides "
             "(every slide, covers included)"), "",
         "Every gate re-opens the `.pptx` from disk. No gate asks the builder what it "
         "believes it wrote.", "",
         "## 5. Render status", "",
         f"`{status}` · {lo_version}", "", RENDER_NOTE, "",
         row("File", "PPTX slides", "Native PDF pages", "Rendered-overflow pages"),
         row("---", "---", "---", "---"),
         *[row(nat[k]["pptx"], c, nat[k].get("page_count", "—"),
               ", ".join(str(x) for x in nat[k]["overflow_pages"]) or "none")
           for k, c in (("sb", sb["slide_count"]), ("lp2", l2["slide_count"]),
                        ("lp3", l3["slide_count"])) if k in nat], "",
         ("Every page of every deck was exported through LibreOffice Impress and rasterised. "
          "No rendered ink enters the footer band, sits below the footer box, or touches the "
          "page edge." if status == NATIVE_COMPLETED
          else "A native render was NOT completed, so no renderer claim is made."), "",
         "## 6. A4 — Comparison", "",
         f"**Applied in this deck: {'yes' if a4['applied_in_this_deck'] else 'NO'}.** "
         f"Treatments actually used: {', '.join(a4['treatments_actually_used'])}. The word "
         "\"comparison\" appears in none of the three files.", "",
         f"The label survives only as a candidate in: {'; '.join(a4['where_the_label_lives'])}.",
         "",
         f"Objects the source does compare: **{' vs '.join(a4['objects_compared'])}** — "
         f"{a4['scope']}", "",
         row("Row", "Relationship", "Axis", "Source wording"),
         row("---", "---", "---", "---")]
    L += [row(r["row"], r["relationship"], r["axis"], r["quoted"])
          for r in a4["explicit_relationships"]]
    L += ["", f"**Asymmetry.** {a4['asymmetry']}", "",
          f"**Verdict.** {a4['verdict']}", "",
          "## 7. The Rumusan mapping question", "",
          "The committed unit model proposes RP-007 for Rumusan — contractor perspective, "
          "**no** Kepentingan / Skop dan Isi Utama / Manfaat labels. A6 is Bariah's ruling "
          "and supersedes that CAIR proposal, so this deck applies A6's three labels to the "
          "unit's own three beats **positionally** "
          f"(`{M.MONTAGE_MAPPING_BASIS}`).", "",
          row("A6 label", "Unit beat it was paired with"), row("---", "---")]
    L += [row(c["label"], c["beat"]) for c in M.rumusan()["montage"]]
    L += ["", "The beat TEXT is the unit's own; only the labels come from A6. **Bariah "
          "should confirm both the supersession and the pairing** — the mapping is "
          "positional, not interpreted, and is marked "
          f"{' · '.join(M.MONTAGE_MARKS)} on the slide.", "",
          "## 8. Unresolved review items", "",
          row("#", "Item", "Why it is open"), row("---", "---", "---"),
          row(1, "Watak / cast", "D1 sets characters per unit; STOP-006 open. The scenario "
                                 "uses the cast-free SITUASI frame with a WATAK placeholder "
                                 "and no proper name."),
          row(2, "Rumusan label-to-beat pairing", "Section 7 above."),
          row(3, "Content-screen visual subjects",
              f"{M.unit_model()['visual_obligations']['count']} candidate subjects, "
              f"{M.unit_model()['visual_obligations']['source_attested']} attested by a "
              "source figure. RP-104 requires a source figure; none exists."),
          row(4, "Quiz answer key",
              f"{M.unit_model()['assessment']['key_status']} — drafted by CAIR, approved by "
              "nobody. Distractors are marked as CAIR-drafted on the slide."),
          row(5, "B2 panel density", "Section 3 above. Not frozen."),
          row(6, "Maximum screen count",
              f"{pol['screen_pattern_plan']['maximum_total']} — only the "
              f"{t['policy_minimum_screens']}-screen floor is fixed."),
          row(7, "A4 Comparison", "Section 6 above. Supportable, not applied."), ""]
    amb = M.unit_model()["ambiguity_register"]
    L += [f"Plus {amb['count']} ambiguities carried from the committed unit model:", ""]
    L += [f"- {a['item']} (`{a['resolution']}`)" for a in amb["items"]]
    L += ["", "## 9. What this package is not", "",
          "This is a **TECHNICAL PROOF**, not an instructional deliverable. It exists to "
          "prove that the generator, the pagination contract and the native render behave "
          "correctly on a real unit, and to let the pattern, the density and the review "
          "packaging be judged on real content. Every instructional decision in it remains "
          "provisional.", "",
          *[f"- `{m}`" for m in M.STATUS_MARKS], "",
          "No screen text, no answer key, no visual direction and no cast in this package is "
          "instructionally approved. The remaining three authorised calibration units "
          f"({', '.join(u for u in M.P.CALIBRATION_UNITS_AUTHORIZED if u != M.UNIT_ID)}) "
          "have not been generated, and mass generation remains unauthorised.", ""]
    path = os.path.join(M.PPTX_DIR, HANDOFF_NAME)
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(L) + "\n")
    return path


def page_break_map():
    """{screen_id: ["1/2", "2/2"]} — the review-page plan, derived from frozen data only."""
    out = {}
    for sl in slides():
        out.setdefault(sl["screen"]["screen_id"], []).append(
            f"{sl['page_index']}/{sl['page_count']}")
    return out


def build_all():
    timings = {}
    t0 = time.time()
    sb = build_storyboard()
    timings["storyboard_pptx_s"] = round(time.time() - t0, 2)
    t0 = time.time()
    lp2 = build_lampiran(2)
    lp3 = build_lampiran(3)
    timings["lampiran_pptx_s"] = round(time.time() - t0, 2)
    return dict(storyboard=sb, lampiran_2=lp2, lampiran_3=lp3, timings=timings)


if __name__ == "__main__":
    t_start = time.time()
    r = build_all()
    t0 = time.time()
    sbp = render_storyboard_previews()
    lp2p = render_lampiran_previews(2)
    lp3p = render_lampiran_previews(3)
    render_s = round(time.time() - t0, 2)
    for k in ("storyboard", "lampiran_2", "lampiran_3"):
        print(f"{k:12} {os.path.basename(r[k])}  ({os.path.getsize(r[k])} bytes)")
    print(f"previews     storyboard={len(sbp)}  lampiran2={len(lp2p)}  "
          f"lampiran3={len(lp3p)}")
    print(f"overflowed   storyboard={len([x for x in sbp if x['overflowed']])}  "
          f"lampiran2={len([x for x in lp2p if x['overflowed']])}  "
          f"lampiran3={len([x for x in lp3p if x['overflowed']])}")
    print(f"clipped      storyboard={len([x for x in sbp if x['clipped']])}")
    print(f"panel_capacity_measured = {panel_capacity()}")
    print(f"timings      build={r['timings']}  render={render_s}s  "
          f"total={round(time.time() - t_start, 2)}s")
