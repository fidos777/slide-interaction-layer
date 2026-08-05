# -*- coding: utf-8 -*-
"""Stage 4.2F-H — QA and mutation fixtures for the K5-PL06-T03-B03 calibration proof.

Every gate re-opens the written PPTX from disk. None of them asks the builder what it
believes it wrote. Where an expectation needs a value rather than a literal, the value comes
from a module that is NOT the calibration model — the raw extraction, the committed unit
model, or Bariah's policy record — so a gate cannot be satisfied by a copy of the thing it is
checking.

Oracle bases used here:

    RAW_FILE            re-read from the .pptx on disk
    DOC_LITERAL         a literal typed into this file by hand
    SOURCE_TEXT_LITERAL the controlled row text from `pl06_extract_v1`
    COMMITTED_MODEL     the committed unit model / policy record, not the calibration model
    INDEPENDENT_MEASURE font metrics, or a recount performed here
"""
import json
import os
import re
import shutil
import sys
import tempfile

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

from pptx import Presentation      # noqa: E402

import k5_calib_model_v1 as M      # noqa: E402
import k5_calib_build_v1 as B      # noqa: E402
import k5_calib_metrics_v1 as MT   # noqa: E402
import pl06_extract_v1 as EX       # noqa: E402
import pl06_unit_model_v1 as U     # noqa: E402
import k5_policy_apply_v1 as AP    # noqa: E402
import k5_pattern_policy_v1 as P   # noqa: E402

SUITE_ID = "K5_CALIBRATION_PROOF_QA_v1"
SIBLING_SUITES = ["K5_BARIAH_POLICY_QA_v1", "T04_STATE_COVERAGE_QA_v1",
                  "STAGE_4_2F_D_RELEASE_FACTS_QA_v1", "STAGE_4_2F_F_K2_QA_v1"]
assert SUITE_ID not in SIBLING_SUITES

RAW_FILE, DOC_LITERAL = "RAW_FILE", "DOC_LITERAL"
SOURCE_TEXT_LITERAL, COMMITTED_MODEL = "SOURCE_TEXT_LITERAL", "COMMITTED_MODEL"
INDEPENDENT_MEASURE = "INDEPENDENT_MEASURE"
NATIVE_RENDER = "NATIVE_RENDER"

# ------------------------------------------------------------------ hand-declared literals
DECLARED_SCREENS = 12
# A learner screen may need more than one REVIEW page. That is documentation pagination and
# does not change the learner-screen count, which stays 12.
DECLARED_REVIEW_PAGES = 17
DECLARED_STORYBOARD_SLIDES = 18          # 1 cover + 17 review pages
DECLARED_MULTIPAGE_SCREENS = ["T03B03-S04", "T03B03-S05", "T03B03-S07", "T03B03-S09",
                              "T03B03-S11"]
DECLARED_RUNTIME_STATES = 43
DECLARED_PANEL_STATES = 38
DECLARED_PRODUCTION_DENSITY = 3
DECLARED_CALIBRATION_DENSITIES = [2, 3]
DECLARED_LAMPIRAN_PAGES = {2: 20, 3: 14}
DECLARED_SLIDE_W, DECLARED_SLIDE_H = 12192000, 6858000
DECLARED_A6_LABELS = ["Kepentingan", "Skop dan Isi Utama", "Manfaat"]
DECLARED_NARRATOR = "Hilmi"
# Hand-typed here on purpose. Reading these off M.STATUS_MARKS / M.MONTAGE_MARKS made the
# gates self-referential: fixtures KR-04 and KA-02 rewrote the constant, the deck followed it,
# and the gates agreed with the mutation. A gate that compares an artifact against its own
# source proves only that a copy succeeded.
DECLARED_STATUS_MARKS = ["TECHNICAL_PROOF", "DRAFT_FOR_BARIAH_REVIEW",
                         "NOT_INSTRUCTIONALLY_APPROVED", "NOT_LEARNER_FACING_FINAL"]
DECLARED_MONTAGE_MARKS = ["PROVISIONAL_VISUAL_DIRECTION", "PENDING_BARIAH_UNIT_REVIEW",
                          "NOT_MMD_ASSET"]
DECLARED_PANEL_FIELDS_MS = ["Pencetus", "Apa yang dipaparkan",
                            "Beza daripada keadaan asas",
                            "Cara pelajar kembali / meneruskan", "Sumber dan kuasa"]
# Names that must never appear: T04's cast, and any cast invented for K5.
FORBIDDEN_NAMES = ["Alya", "Encik Rahman", "Puan ", "Cikgu ", "Ahmad", "Siti"]
FORBIDDEN_T04_TOKENS = ["T04-", "T04_", "K5PL06T04"]
FORBIDDEN_FINALITY = ["Tamat Topik", "DILULUSKAN OLEH BARIAH", "KUNCI JAWAPAN DISAHKAN"]
DECLARED_TREATMENTS = ["click-to-reveal", "scenario/dialogue only where justified",
                       "static content"]
DECLARED_A4_ROWS = ["T03B03-ROW-056", "T03B03-ROW-061"]
DECLARED_A4_WORDING = {"T03B03-ROW-056": "alternatif mesra alam kepada parit konkrit",
                       "T03B03-ROW-061": "bukannya parit yang dalam dan curam"}
HANDOFF_NAME = "K5_PL06_T03_B03_CALIBRATION_HANDOFF_v0_1.md"

# ---------------------------------------------------------------- native-oracle thresholds
# Hand-typed HERE, in points, and passed INTO the render check. They were previously read off
# k5_calib_build_v1's own geometry constants — so if the generator moved the footer, the
# oracle moved with it and could never disagree. An oracle that imports its thresholds from
# the module under test is not an oracle.
#   storyboard body box   1,320,000 + 4,600,000 EMU  -> 466.14 pt
#   Lampiran panel box    1,280,000 + 4,700,000 EMU  -> 470.87 pt
#   footer box            6,150,000 + 500,000 EMU    -> 484.25 pt .. 523.62 pt
DECLARED_BANDS_PT = {
    "sb": (466.14, 484.25, 523.62),
    "lp2": (470.87, 484.25, 523.62),
    "lp3": (470.87, 484.25, 523.62),
}
DECLARED_PAGE_H_PT = 540.0
# Every bounded text body in all three decks must declare noAutofit. There are no exemptions;
# if one is ever needed it is named here and in the handoff, never left implicit.
AUTOFIT_EXCEPTIONS = []
DECLARED_LINE_FACTOR = 1.2
DECLARED_METRICS_VERSION = "1"
DECLARED_FONT_SHA256 = ("4659bc0c58c5028dd488ec928d41d9265db43d9b669fc14ca8b0832"
                        "daca7b144")


# ==========================================================================================
class Gate(dict):
    pass


def _g(name, gate_type, observed, expected, population, source_artifact, oracle_basis,
       empty_by_design=False):
    ok = observed == expected
    if population == 0 and not empty_by_design:
        ok = False
    return Gate(name=name, gate_type=gate_type, population=population,
                expected=expected, observed=observed, source_artifact=source_artifact,
                oracle_basis=oracle_basis, empty_by_design=empty_by_design, passed=ok)


def _paths():
    return dict(
        sb=os.path.join(M.PPTX_DIR, M.STORYBOARD_NAME),
        lp2=os.path.join(M.PPTX_DIR, M.LAMPIRAN_NAME.format(n=2)),
        lp3=os.path.join(M.PPTX_DIR, M.LAMPIRAN_NAME.format(n=3)))


_EX_ONCE = None
_UM_ONCE = None


def _inputs():
    """The committed inputs, parsed once. They are invariant across mutation fixtures."""
    global _EX_ONCE, _UM_ONCE
    if _EX_ONCE is None:
        _EX_ONCE = EX.extract_all()[M.UNIT_ID]
        _UM_ONCE = U.model(M.UNIT_ID, _EX_ONCE)
    return _EX_ONCE, _UM_ONCE


def _expected_screen_ids(pol):
    """Screen IDs derived from the POLICY record and the unit ID, not from the screen roll.

    shell(after D3/A5) + content groups + closing, numbered from the unit's short code.
    """
    plan = pol["screen_pattern_plan"]
    n = (plan["shell_screens_after"] + plan["content_groups"] + plan["closing_screens"])
    short = "".join(M.UNIT_ID.split("-")[2:])
    return [f"{short}-S{i:02d}" for i in range(1, n + 1)]


PARITY_NORMALISATION = ("A soft line break serialises as a vertical tab in the run text; it "
                        "is normalised to a space on readback. Nothing else is normalised.")


def _parity_report():
    """Replay a traced build into scratch and compare canonical / measured / rendered / read.

    Never writes to the real package directory.
    """
    import collections
    tmp = tempfile.mkdtemp(prefix="k5parity_")
    old_dir, old_prev = M.PPTX_DIR, B.WRITE_PREVIEW_IMAGES
    M.PPTX_DIR, B.WRITE_PREVIEW_IMAGES = tmp, False
    try:
        B.trace_start()
        path = B.build_storyboard()
        tr = B.trace_stop()
        read = [t for _, t in B.storyboard_body_runs(path)]
    finally:
        M.PPTX_DIR, B.WRITE_PREVIEW_IMAGES = old_dir, old_prev
        shutil.rmtree(tmp, ignore_errors=True)

    blocks = [b for sc in M.screens() for b in sc["blocks"]]
    labels = tr["label"]
    ids = [labels.get(id(b), f"UNLABELLED:{i}") for i, b in enumerate(blocks)]
    dup = [k for k, v in collections.Counter(ids).items() if v > 1]
    missing = [labels.get(id(b), f"UNLABELLED:{i}") for i, b in enumerate(blocks)
               if id(b) not in tr["measured"] or id(b) not in tr["rendered"]]

    remaining = collections.Counter(read)
    mismatched = []
    for i, b in enumerate(blocks):
        lab = labels.get(id(b), f"UNLABELLED:{i}")
        canon = B.block_text(b)
        meas = tr["measured"].get(id(b))
        rend = tr["rendered"].get(id(b))
        in_file = remaining.get(rend, 0) > 0 if rend is not None else False
        if in_file:
            remaining[rend] -= 1
        if not (canon == meas == rend) or not in_file:
            mismatched.append(dict(block=lab, canonical=canon, measured=meas,
                                   rendered=rend, found_in_pptx=in_file))
    return dict(population=len(blocks), missing=missing, duplicate=dup,
                mismatched=mismatched, readback_runs=len(read))


def _screen_of(kind):
    """The screen of a kind, or None. A gate must be able to report an ABSENT screen rather
    than raise before it is evaluated — fixture KN-04 deletes the scenario screen."""
    return next((s for s in M.screens() if s["kind"] == kind), None)


NATIVE_DIR = os.path.join(tempfile.gettempdir(), "k5_native_render")


def gates(native=None):
    p = _paths()
    rb = {k: B.readback(v) for k, v in p.items()}
    sb, lp2, lp3 = rb["sb"], rb["lp2"], rb["lp3"]
    all_text = {k: v["all_text"] for k, v in rb.items()}
    ex, um = _inputs()
    pol = next(u for u in AP.calibration_units()["units"] if u["unit_id"] == M.UNIT_ID)
    G = []
    native_ran = native is not None

    # ---------------------------------------------------------------- files and geometry ---
    G.append(_g("ALL_THREE_FILES_EXIST_AND_ARE_NON_EMPTY", "READBACK",
                sorted(k for k, v in p.items() if os.path.exists(v)
                       and os.path.getsize(v) > 0),
                ["lp2", "lp3", "sb"], 3, "the three .pptx on disk", RAW_FILE))
    G.append(_g("EVERY_FILE_IS_16_BY_9_AT_THE_DECLARED_EMU_SIZE", "READBACK",
                sorted({(v["slide_width"], v["slide_height"]) for v in rb.values()}),
                [(DECLARED_SLIDE_W, DECLARED_SLIDE_H)], 3, "readback", DOC_LITERAL))
    G.append(_g("STORYBOARD_SLIDE_COUNT_MATCHES_THE_DECLARED_ROLL", "READBACK",
                sb["slide_count"], DECLARED_STORYBOARD_SLIDES, 1, M.STORYBOARD_NAME,
                DOC_LITERAL))
    for n, f in ((2, lp2), (3, lp3)):
        G.append(_g(f"LAMPIRAN_{n}_PANEL_SLIDE_COUNT_MATCHES_THE_DECLARED_ROLL", "READBACK",
                    f["slide_count"], DECLARED_LAMPIRAN_PAGES[n] + 1, 1,
                    M.LAMPIRAN_NAME.format(n=n), DOC_LITERAL))
    G.append(_g("NO_SHAPE_SITS_OFF_CANVAS_IN_ANY_FILE", "READBACK",
                sum(len(v["off_canvas"]) for v in rb.values()), 0, 3, "readback geometry",
                RAW_FILE, empty_by_design=False))
    G.append(_g("EVERY_STORYBOARD_SCREEN_SLIDE_CARRIES_SPEAKER_NOTES", "READBACK",
                len([s for s in sb["slides"][1:] if s["notes"].strip()]),
                DECLARED_STORYBOARD_SLIDES - 1, DECLARED_STORYBOARD_SLIDES - 1,
                M.STORYBOARD_NAME, RAW_FILE))
    G.append(_g("SPEAKER_NOTES_CARRY_THE_VO_AND_THE_SOURCE_ROWS", "READBACK",
                len([s for s in sb["slides"][1:]
                     if "VO (" in s["notes"] and "Baris sumber:" in s["notes"]]),
                DECLARED_STORYBOARD_SLIDES - 1, DECLARED_STORYBOARD_SLIDES - 1,
                M.STORYBOARD_NAME, RAW_FILE))

    # ------------------------------------------------------------- source traceability -----
    rows = [r for r in ex["rows"] if (r["controlled_display_text"] or "").strip()]
    missing = [r["row_id"] for r in rows
               if r["controlled_display_text"].strip() not in all_text["sb"]]
    G.append(_g("EVERY_CONTROLLED_SOURCE_ROW_APPEARS_ON_THE_STORYBOARD", "TRACEABILITY",
                missing, [], len(rows), M.STORYBOARD_NAME, SOURCE_TEXT_LITERAL))
    by = {r["row_id"]: (r["controlled_display_text"] or "").strip() for r in ex["rows"]}
    verbatim = [b for sc in M.screens() for b in sc["blocks"]
                if b["provenance"] == "SOURCE_CONTROLLED"]
    bad = [b.get("row_id") or f"NO_ROW_ID:{b['text'][:40]}"
           for b in verbatim if by.get(b.get("row_id")) != b["text"]]
    G.append(_g("EVERY_VERBATIM_BLOCK_EQUALS_ITS_CONTROLLED_ROW", "TRACEABILITY",
                bad, [], len(verbatim), "block roll vs pl06_extract_v1",
                SOURCE_TEXT_LITERAL))
    cond = [b for sc in M.screens() for b in sc["blocks"]
            if b["provenance"] == "SOURCE_CONDENSED"]
    added = [b["text"] for b in cond
             if not b.get("condensed_from", "").startswith(b["text"])]
    G.append(_g("NO_CONDENSED_LINE_ADDS_A_WORD_TO_ITS_PROPOSITION", "TRACEABILITY",
                added, [], len(cond), "block roll vs committed propositions",
                COMMITTED_MODEL))
    cited = {r for sc in M.screens() for r in sc["source_row_ids"]}
    G.append(_g("EVERY_CITED_ROW_ID_IS_A_REAL_ROW_OF_THIS_UNIT", "TRACEABILITY",
                sorted(cited - set(by)), [], len(cited), "citations vs pl06_extract_v1",
                SOURCE_TEXT_LITERAL))
    # S01 cited a row in a block while declaring no source rows, so its Notes read
    # "skrin struktur" over text that came straight from the module.
    undeclared = [(sc["screen_id"], b["row_id"]) for sc in M.screens()
                  for b in sc["blocks"]
                  if b.get("row_id") and b["row_id"] not in sc["source_row_ids"]]
    G.append(_g("EVERY_ROW_A_BLOCK_CITES_IS_DECLARED_BY_ITS_SCREEN", "TRACEABILITY",
                undeclared, [], len([b for sc in M.screens() for b in sc["blocks"]
                                     if b.get("row_id")]),
                "block roll vs screen roll", INDEPENDENT_MEASURE))
    G.append(_g("NO_BLOCK_CARRIES_THE_FORBIDDEN_PROVENANCE", "TRACEABILITY",
                len([b for sc in M.screens() for b in sc["blocks"]
                     if b["provenance"] == M.FORBIDDEN_PROVENANCE]), 0,
                len([b for sc in M.screens() for b in sc["blocks"]]),
                "block roll", DOC_LITERAL))
    G.append(_g("EVERY_BLOCK_USES_A_DECLARED_PROVENANCE_CLASS", "TRACEABILITY",
                sorted({b["provenance"] for sc in M.screens() for b in sc["blocks"]}
                       - set(M.PROVENANCE)), [],
                len([b for sc in M.screens() for b in sc["blocks"]]),
                "block roll", DOC_LITERAL))

    # ------------------------------------------------------------------ state coverage -----
    sts = M.states()
    G.append(_g("RUNTIME_STATE_COUNT_MATCHES_THE_DECLARED_ROLL", "STATE",
                len(sts), DECLARED_RUNTIME_STATES, len(sts), "k5 state roll", DOC_LITERAL))
    counts = {}
    for s in sts:
        counts[s["screen_id"]] = counts.get(s["screen_id"], 0) + 1
    recount = len([s for s in sts if counts[s["screen_id"]] > 1])
    G.append(_g("PANEL_STATE_COUNT_RECOUNTS_TO_THE_DECLARED_ROLL", "STATE",
                (len(M.panel_states()), recount),
                (DECLARED_PANEL_STATES, DECLARED_PANEL_STATES), len(sts),
                "state roll recounted here", INDEPENDENT_MEASURE))
    for n, f in ((2, lp2), (3, lp3)):
        absent = [s["state_id"] for s in M.panel_states()
                  if s["state_id"] not in f["all_text"]]
        G.append(_g(f"EVERY_PANEL_STATE_APPEARS_IN_THE_{n}_PANEL_LAMPIRAN", "STATE",
                    absent, [], DECLARED_PANEL_STATES, M.LAMPIRAN_NAME.format(n=n),
                    RAW_FILE))
    ids2 = {t for t in lp2["all_text"].split("\n") if "-ST-" in t}
    ids3 = {t for t in lp3["all_text"].split("\n") if "-ST-" in t}
    G.append(_g("BOTH_DENSITIES_CARRY_THE_SAME_STATES_ONLY_THE_LAYOUT_DIFFERS", "STATE",
                sorted(ids2 ^ ids3), [], len(ids2 | ids3),
                "2-panel vs 3-panel readback", RAW_FILE))
    G.append(_g("THE_TWO_DENSITIES_DIFFER_IN_PAGE_COUNT", "STATE",
                lp2["slide_count"] != lp3["slide_count"], True, 2, "readback", RAW_FILE))
    for n in (2, 3):
        pages = M.lampiran_pages(n)
        mixed = [pg["page_no"] for pg in pages
                 if len({s["screen_id"] for s in pg["states"]}) > 1]
        over = [pg["page_no"] for pg in pages if len(pg["states"]) > n]
        G.append(_g(f"NO_{n}_PANEL_PAGE_MIXES_TWO_SCREENS_OR_EXCEEDS_ITS_DENSITY", "STATE",
                    (mixed, over), ([], []), len(pages),
                    M.LAMPIRAN_NAME.format(n=n), INDEPENDENT_MEASURE))
    for n, f in ((2, lp2), (3, lp3)):
        short = [s["index"] for s in f["slides"][1:]
                 if not all(lbl in "\n".join(s["texts"])
                            for lbl in DECLARED_PANEL_FIELDS_MS)]
        G.append(_g(f"EVERY_{n}_PANEL_PAGE_ANSWERS_ALL_FIVE_COVERAGE_FIELDS", "STATE",
                    short, [], DECLARED_LAMPIRAN_PAGES[n],
                    M.LAMPIRAN_NAME.format(n=n), RAW_FILE))

    # ----------------------------------------------------------------------- scenario ------
    sc2 = _screen_of("SCENARIO")
    G.append(_g("A_SCENARIO_SCREEN_EXISTS_AS_D3_REQUIRES", "SCENARIO",
                (sc2["screen_id"], "D3" in sc2["decision_ids"]) if sc2
                else ("ABSENT", False),
                ("T03B03-S02", True), 1, "screen roll", DOC_LITERAL))
    G.append(_g("THE_SCENARIO_CARRIES_A_PENDING_CAST_PLACEHOLDER_ON_THE_FACE_OF_THE_DECK",
                "SCENARIO",
                "WATAK: PENDING_UNIT_REVIEW" in all_text["sb"], True, 1,
                M.STORYBOARD_NAME, RAW_FILE))
    hits = [n for n in FORBIDDEN_NAMES
            if any(n in t for t in all_text.values())]
    G.append(_g("NO_PROPER_NAME_WAS_INVENTED_FOR_THE_CAST", "SCENARIO",
                hits, [], len(FORBIDDEN_NAMES), "all three files", DOC_LITERAL))
    sitems = sc2["scenario"]["items"] if sc2 else []
    unsourced = [i["row_id"] for i in sitems if by.get(i["row_id"]) != i["text"]]
    G.append(_g("EVERY_SCENARIO_LINE_IS_A_CONTROLLED_ROW_VERBATIM", "SCENARIO",
                unsourced, [], len(sitems),
                "scenario vs pl06_extract_v1", SOURCE_TEXT_LITERAL))
    G.append(_g("THE_SCENARIO_IS_MARKED_DRAFT_ON_THE_DECK", "SCENARIO",
                all(m in all_text["sb"] for m in DECLARED_STATUS_MARKS), True,
                len(DECLARED_STATUS_MARKS), M.STORYBOARD_NAME, DOC_LITERAL))

    # ------------------------------------------------------------------------ Rumusan ------
    _rs = _screen_of("RUMUSAN")
    rm = _rs["rumusan"] if _rs else dict(montage=[], beats=[])
    G.append(_g("THE_MONTAGE_USES_A6_THREE_LABELS_IN_ORDER", "RUMUSAN",
                [c["label"] for c in rm["montage"]], DECLARED_A6_LABELS, 3,
                "montage roll", DOC_LITERAL))
    committed = um["rumusan"]["beats"]
    G.append(_g("EVERY_MONTAGE_COMPONENT_POINTS_AT_A_COMMITTED_UNIT_BEAT", "RUMUSAN",
                [c["beat"] for c in rm["montage"] if c["beat"] not in committed], [],
                len(rm["montage"]), "montage vs pl06_unit_model_v1", COMMITTED_MODEL))
    G.append(_g("THE_MONTAGE_INVENTS_NO_FOURTH_SUBJECT", "RUMUSAN",
                len(rm["montage"]), len(committed), len(committed),
                "montage vs pl06_unit_model_v1", COMMITTED_MODEL))
    G.append(_g("THE_MONTAGE_IS_MARKED_PROVISIONAL_ON_THE_DECK", "RUMUSAN",
                [m for m in DECLARED_MONTAGE_MARKS if m not in all_text["sb"]], [],
                len(DECLARED_MONTAGE_MARKS), M.STORYBOARD_NAME, DOC_LITERAL))
    G.append(_g("THE_MODELS_MONTAGE_MARKS_HAVE_NOT_DRIFTED_FROM_THE_LITERALS", "RUMUSAN",
                list(M.MONTAGE_MARKS), DECLARED_MONTAGE_MARKS, len(DECLARED_MONTAGE_MARKS),
                "k5_calib_model_v1.MONTAGE_MARKS", DOC_LITERAL))

    # ----------------------------------------------------------------------- narrator ------
    G.append(_g("THE_NARRATOR_ON_THE_DECK_IS_THE_ONE_D2_NAMES", "NARRATOR",
                (pol["narrator"]["name"], f"Narator: {pol['narrator']['name']}"
                 in all_text["sb"]),
                (DECLARED_NARRATOR, True), 1, M.STORYBOARD_NAME, DOC_LITERAL))
    notes_with_vo = [s for s in sb["slides"][1:]
                     if f"VO ({DECLARED_NARRATOR})" in s["notes"]]
    G.append(_g("EVERY_SCREEN_ATTRIBUTES_ITS_VO_TO_THE_NAMED_NARRATOR", "NARRATOR",
                len(notes_with_vo), DECLARED_STORYBOARD_SLIDES - 1,
                DECLARED_STORYBOARD_SLIDES - 1, M.STORYBOARD_NAME, RAW_FILE))

    # --------------------------------------------------------------------------- quiz ------
    items = um["assessment"]["items"]
    G.append(_g("THE_QUIZ_IS_FOUR_MCQ_PLUS_ONE_MULTIPLE_RESPONSE", "QUIZ",
                (len([q for q in items if q["kind"] == "MULTIPLE_CHOICE"]),
                 len([q for q in items if q["kind"] == "MULTIPLE_RESPONSE"])),
                (4, 1), len(items), "pl06_unit_model_v1 assessment", COMMITTED_MODEL))
    G.append(_g("EVERY_COMMITTED_STEM_APPEARS_ON_THE_DECK", "QUIZ",
                [q["slot"] for q in items if q["stem"] not in all_text["sb"]], [],
                len(items), M.STORYBOARD_NAME, COMMITTED_MODEL))
    G.append(_g("THE_ANSWER_KEY_IS_SHOWN_AS_DRAFTED_NOT_APPROVED", "QUIZ",
                f"KUNCI JAWAPAN: {um['assessment']['key_status']}" in all_text["sb"]
                and um["assessment"]["key_status"] == "DRAFTED_NOT_APPROVED",
                True, 1, M.STORYBOARD_NAME, RAW_FILE))
    G.append(_g("THE_PASS_MARK_ON_THE_DECK_IS_SIXTY_PERCENT", "QUIZ",
                "Markah lulus 60 peratus" in all_text["sb"], True, 1,
                M.STORYBOARD_NAME, DOC_LITERAL))
    G.append(_g("EVERY_DISTRACTOR_IS_MARKED_AS_CAIR_DRAFTED", "QUIZ",
                sorted({b["provenance"] for sc in M.screens() if sc["kind"] == "KUIZ"
                        for b in sc["blocks"]
                        if b["text"] in [d for q in items for d in q["distractors"]]}),
                ["CAIR_DRAFTED_ASSESSMENT"],
                len([d for q in items for d in q["distractors"]]),
                "block roll", DOC_LITERAL))

    # ------------------------------------------------------------------------ closing ------
    cl = pol["closing_screen"]
    _tm = _screen_of("TAMAT")
    G.append(_g("THE_CLOSING_TEXT_IS_D4S_INTERMEDIATE_BAHAGIAN_STRING", "CLOSING",
                (bool(_tm) and cl["closing_text"] in all_text["sb"],
                 cl["is_final_bahagian"]),
                (True, False), 1, "deck vs k5_policy_apply_v1", COMMITTED_MODEL))
    G.append(_g("THE_DECK_DOES_NOT_CLOSE_THE_TOPIK", "CLOSING",
                [s for s in FORBIDDEN_FINALITY if any(s in t for t in all_text.values())],
                [], len(FORBIDDEN_FINALITY), "all three files", DOC_LITERAL))
    G.append(_g("THE_BAHAGIAN_LABEL_IS_SHOWN_BECAUSE_THE_TOPIK_HAS_MANY", "CLOSING",
                (cl["bahagian_label_shown"], cl["segmentation"]),
                (True, "TOPIK_WITH_MULTIPLE_BAHAGIAN"), 1, "k5_policy_apply_v1",
                COMMITTED_MODEL))

    # ---------------------------------------------------------- no approval is claimed -----
    # A plain substring search matches the negated marks: NOT_INSTRUCTIONALLY_APPROVED
    # contains INSTRUCTIONALLY_APPROVED. The claim only counts when it is not negated.
    claims = [c for c in M.FORBIDDEN_STATUS_CLAIMS
              if any(re.search(r"(?<!NOT_)" + re.escape(c), t) for t in all_text.values())]
    G.append(_g("NO_FILE_CLAIMS_ANY_FORM_OF_APPROVAL", "STATUS",
                claims, [], len(M.FORBIDDEN_STATUS_CLAIMS), "all three files", DOC_LITERAL))
    G.append(_g("EVERY_FILE_CARRIES_EVERY_DECLARED_STATUS_MARK", "STATUS",
                sorted(k for k, t in all_text.items()
                       if all(m in t for m in DECLARED_STATUS_MARKS)),
                ["lp2", "lp3", "sb"], 3, "all three files", DOC_LITERAL))
    G.append(_g("THE_MODELS_STATUS_MARKS_HAVE_NOT_DRIFTED_FROM_THE_LITERALS", "STATUS",
                list(M.STATUS_MARKS), DECLARED_STATUS_MARKS, len(DECLARED_STATUS_MARKS),
                "k5_calib_model_v1.STATUS_MARKS", DOC_LITERAL))
    # B2 asked WHICH density to ship and was answered at F5 of the returned authority:
    # "3 panel sebagai format produksi." The Kelompok 0 record still reads TEST_REQUIRED,
    # correctly — it is a 4 August record and F5 is 5 August. The gate now checks that the
    # historical record is unchanged AND that the production density comes from F5.
    G.append(_g("B2_HISTORICAL_RECORD_IS_UNCHANGED_AND_F5_SUPPLIES_THE_PRODUCTION_DENSITY",
                "STATUS",
                (next(d for d in P.DECISIONS if d["id"] == "B2")["status"],
                 B.PRODUCTION_PANEL_DENSITY, sorted(B.CALIBRATION_DENSITIES)),
                ("TEST_REQUIRED_BEFORE_FINAL_FREEZE", DECLARED_PRODUCTION_DENSITY,
                 DECLARED_CALIBRATION_DENSITIES), 1,
                "k5_pattern_policy_v1 + the authority declaration", COMMITTED_MODEL))
    G.append(_g("THE_TWO_PANEL_FILE_IS_RETAINED_AS_CALIBRATION_EVIDENCE", "STATUS",
                os.path.exists(os.path.join(M.PPTX_DIR, M.LAMPIRAN_NAME.format(n=2))),
                True, 1, "on-disk bytes", RAW_FILE))
    G.append(_g("MASS_GENERATION_REMAINS_UNAUTHORISED", "STATUS",
                (P.CALIBRATION_GENERATION_AUTHORIZED, P.MASS_GENERATION_AUTHORIZED),
                (True, False), 1, "k5_pattern_policy_v1", COMMITTED_MODEL))

    # ------------------------------------------------------------- nothing T04 leaked ------
    leaked = [t for t in FORBIDDEN_T04_TOKENS if any(t in x for x in all_text.values())]
    G.append(_g("NO_T04_IDENTIFIER_APPEARS_IN_ANY_K5_FILE", "NON_INHERITANCE",
                leaked, [], len(FORBIDDEN_T04_TOKENS), "all three files", DOC_LITERAL))
    G.append(_g("THE_SCREEN_COUNT_COMES_FROM_BARIAHS_POLICY_NOT_FROM_T04", "NON_INHERITANCE",
                (len(M.screens()), pol["screen_pattern_plan"]["minimum_total_after"]),
                (DECLARED_SCREENS, DECLARED_SCREENS), DECLARED_SCREENS,
                "screen roll vs k5_policy_apply_v1", COMMITTED_MODEL))
    G.append(_g("THE_STATE_COUNT_IS_NOT_T04S_SIXTY_FIVE", "NON_INHERITANCE",
                len(sts) == 65, False, len(sts), "state roll", DOC_LITERAL))
    G.append(_g("THE_CONTENT_GROUP_COUNT_IS_THIS_UNITS_OWN", "NON_INHERITANCE",
                (len(M.groups()), len(um["grouping"]["groups"])), (6, 6), 6,
                "groups vs pl06_unit_model_v1", COMMITTED_MODEL))

    # ----------------------------------------------------------------------------- A4 ------
    G.append(_g("NO_SCREEN_ASSERTS_A_COMPARISON_TREATMENT", "A4",
                sorted({s["treatment"] for s in M.screens()}), DECLARED_TREATMENTS,
                len(M.screens()), "screen roll", DOC_LITERAL))
    G.append(_g("THE_WORD_COMPARISON_APPEARS_IN_NO_GENERATED_FILE", "A4",
                sorted(k for k, t in all_text.items()
                       if "comparison" in t.lower() or "perbandingan" in t.lower()),
                [], 3, "all three files", RAW_FILE, empty_by_design=True))
    G.append(_g("THE_MODEL_DOES_NOT_CLAIM_A4_WAS_APPLIED", "A4",
                M.A4_COMPARISON["applied_in_this_deck"], False, 1,
                "k5_calib_model_v1.A4_COMPARISON", DOC_LITERAL))
    G.append(_g("EVERY_A4_EVIDENCE_ROW_EXISTS_AND_CARRIES_ITS_COMPARATIVE_WORDING", "A4",
                [r for r in DECLARED_A4_ROWS
                 if DECLARED_A4_WORDING[r] not in by.get(r, "")], [],
                len(DECLARED_A4_ROWS), "A4 rows vs pl06_extract_v1", SOURCE_TEXT_LITERAL))
    G.append(_g("THE_A4_EVIDENCE_ROWS_ARE_THE_ONES_THE_MODEL_CITES", "A4",
                sorted(M.A4_COMPARISON["supporting_rows"]), sorted(DECLARED_A4_ROWS),
                len(DECLARED_A4_ROWS), "k5_calib_model_v1.A4_COMPARISON", DOC_LITERAL))
    # The compared object is described only by negation inside the swale's own rows. If a row
    # ever describes a concrete drain in its own right, the comparison changes character and
    # this gate should be revisited rather than silently kept passing.
    independent = [r["row_id"] for r in ex["rows"]
                   if re.search(r"parit konkrit|longkang konkrit",
                                r["controlled_display_text"] or "", re.I)
                   and r["row_id"] not in DECLARED_A4_ROWS]
    G.append(_g("THE_COMPARED_OBJECT_HAS_NO_INDEPENDENT_SOURCE_DESCRIPTION", "A4",
                independent, [], len(ex["rows"]), "pl06_extract_v1", SOURCE_TEXT_LITERAL))

    # ------------------------------------------------------------------------ handoff ------
    hp = os.path.join(M.PPTX_DIR, HANDOFF_NAME)
    htxt = open(hp, encoding="utf-8").read() if os.path.exists(hp) else ""
    required = [M.STORYBOARD_NAME, M.LAMPIRAN_NAME.format(n=2), M.LAMPIRAN_NAME.format(n=3),
                "RENDER_STATUS", "Native PDF pages", "Rendered-overflow pages",
                "Storyboard review pages", "Learner screens", "A4",
                "Rumusan", "DRAFT_FOR_BARIAH_REVIEW"]
    G.append(_g("THE_HANDOFF_LISTS_EVERY_REQUIRED_ITEM", "HANDOFF",
                [r for r in required if r not in htxt], [], len(required),
                HANDOFF_NAME, DOC_LITERAL))
    G.append(_g("THE_HANDOFF_CLAIMS_NO_NATIVE_POWERPOINT_VALIDATION", "HANDOFF",
                [c for c in M.FORBIDDEN_STATUS_CLAIMS
                 if re.search(r"(?<!NOT_)" + re.escape(c), htxt)]
                + (["CLAIMS_NATIVE_POWERPOINT"] if "rendered in PowerPoint" in htxt else []),
                [], len(M.FORBIDDEN_STATUS_CLAIMS), HANDOFF_NAME, DOC_LITERAL))

    # ------------------------------------------------------------------------- layout ------
    sbp = B.render_storyboard_previews()
    lp2p = B.render_lampiran_previews(2)
    lp3p = B.render_lampiran_previews(3)
    G.append(_g("NO_STORYBOARD_PAGE_OVERFLOWS_ITS_COLUMN", "LAYOUT",
                [x["screen_id"] for x in sbp if x["overflowed"]], [], len(sbp),
                "re-measured preview", INDEPENDENT_MEASURE))
    G.append(_g("NO_STORYBOARD_LINE_IS_CLIPPED_AT_THE_COLUMN_EDGE", "LAYOUT",
                [x["screen_id"] for x in sbp if x["clipped"]], [], len(sbp),
                "re-measured preview", INDEPENDENT_MEASURE))
    for n, pv in ((2, lp2p), (3, lp3p)):
        G.append(_g(f"NO_{n}_PANEL_PAGE_OVERFLOWS_ITS_PANEL_BOX", "LAYOUT",
                    [x["page_no"] for x in pv if x["overflowed"]], [], len(pv),
                    "re-measured preview", INDEPENDENT_MEASURE))
    cap = B.panel_capacity()
    G.append(_g("BOTH_TESTED_DENSITIES_SIT_INSIDE_THE_MEASURED_RENDERER_CAPACITY", "LAYOUT",
                (cap >= 3, cap), (True, cap), 2,
                "t04_state_emit_v1._panel_wrapped_height on K5 states",
                INDEPENDENT_MEASURE))
    # Preview coverage. The count on the left comes from the file on disk, the count on the
    # right from the renderer's own returned records — neither is the renderer's belief about
    # the other. Two Lampiran covers were silently unrendered before this gate existed.
    cover = [(sb["slide_count"], len(sbp)), (lp2["slide_count"], len(lp2p)),
             (lp3["slide_count"], len(lp3p))]
    G.append(_g("EVERY_PPTX_SLIDE_HAS_A_RENDERED_PREVIEW", "LAYOUT",
                [(a, b) for a, b in cover if a != b], [],
                sum(a for a, _ in cover), "readback vs render records",
                INDEPENDENT_MEASURE))
    on_disk = [x["path"] for x in sbp + lp2p + lp3p if not os.path.exists(x["path"])]
    G.append(_g("EVERY_RENDERED_PREVIEW_WAS_WRITTEN_TO_DISK", "LAYOUT",
                on_disk if B.WRITE_PREVIEW_IMAGES else [], [],
                len(sbp + lp2p + lp3p) if B.WRITE_PREVIEW_IMAGES else 0,
                "preview directories", RAW_FILE,
                empty_by_design=not B.WRITE_PREVIEW_IMAGES))
    # ---------------------------------------------------------------------- structural ----
    # POSITIVE contract. The old gate only forbade spAutoFit, so a bodyPr with no autofit
    # child at all passed while declaring nothing. The package contract is that every bounded
    # text body states `noAutofit` explicitly; omission is ambiguous and therefore fails. No
    # claim is made about what a renderer infers from an empty bodyPr.
    audit = {k: B.autofit_audit(v) for k, v in p.items()}
    bodies = [(k, r) for k, rs in audit.items() for r in rs]
    bad = [f"{k}:s{r['slide']}:{r['shape_id']}"
           for k, r in bodies
           if not (r["noAutofit"] and not r["spAutoFit"] and not r["normAutofit"])]
    G.append(_g("EVERY_BOUNDED_TEXT_BOX_EXPLICITLY_DECLARES_NO_AUTOFIT", "STRUCTURAL",
                bad, [], len(bodies), "bodyPr read back from the .pptx", RAW_FILE))
    G.append(_g("THE_AUTOFIT_CONTRACT_HAS_NO_PERMITTED_EXCEPTIONS", "STRUCTURAL",
                AUTOFIT_EXCEPTIONS, [], len(bodies), "k5_calib_qa_v1", DOC_LITERAL,
                empty_by_design=True))
    G.append(_g("LEARNER_SCREEN_COUNT_IS_UNCHANGED_BY_REVIEW_PAGINATION", "STRUCTURAL",
                (len(M.screens()), len(B.slides())),
                (DECLARED_SCREENS, DECLARED_REVIEW_PAGES), DECLARED_SCREENS,
                "screen roll vs review-page roll", DOC_LITERAL))
    labelled = [sl["screen"]["screen_id"] for sl in B.slides()
                if sl["page_count"] > 1 and "halaman" not in sl["title_suffix"]]
    G.append(_g("EVERY_MULTIPAGE_SCREEN_LABELS_ITS_REVIEW_PAGE", "STRUCTURAL",
                labelled, [],
                len([sl for sl in B.slides() if sl["page_count"] > 1]),
                "review-page roll", INDEPENDENT_MEASURE))
    G.append(_g("THE_MULTIPAGE_SCREENS_ARE_THE_DECLARED_ONES", "STRUCTURAL",
                sorted({sl["screen"]["screen_id"] for sl in B.slides()
                        if sl["page_count"] > 1}),
                DECLARED_MULTIPAGE_SCREENS, len(DECLARED_MULTIPAGE_SCREENS),
                "review-page roll", DOC_LITERAL))

    # ------------------------------------------------------------------- native render ----
    # The only oracle here is what LibreOffice actually drew. No model geometry, no preview.
    nat = native or {}
    G.append(_g("THE_NATIVE_RENDER_COMPLETED_FOR_EVERY_FILE", "NATIVE",
                sorted(k for k, r in nat.items() if r["status"] != B.NATIVE_COMPLETED),
                [], len(nat), "LibreOffice Impress export", NATIVE_RENDER,
                empty_by_design=not native_ran))
    G.append(_g("EVERY_NATIVE_PDF_PAGE_COUNT_EQUALS_ITS_PPTX_SLIDE_COUNT", "NATIVE",
                sorted(k for k in nat if nat[k].get("page_count") != rb[k]["slide_count"]),
                [], len(nat), "PDF page count vs readback", NATIVE_RENDER,
                empty_by_design=not native_ran))
    G.append(_g("NO_RENDERED_INK_ENTERS_THE_FOOTER_BAND", "NATIVE",
                sorted((k, f["page"]) for k, r in nat.items() for f in r["findings"]
                       if f["ink_in_gap_band"] or f["ink_below_footer"]),
                [], sum(len(r["findings"]) for r in nat.values()),
                "rasterised PDF pages", NATIVE_RENDER,
                empty_by_design=not native_ran))
    G.append(_g("NO_RENDERED_INK_LEAVES_THE_PAGE_CANVAS", "NATIVE",
                sorted((k, f["page"]) for k, r in nat.items() for f in r["findings"]
                       if f["ink_at_page_edge"]),
                [], sum(len(r["findings"]) for r in nat.values()),
                "rasterised PDF pages", NATIVE_RENDER,
                empty_by_design=not native_ran))
    # Silent clipping is the risk noAutofit introduces: the box no longer grows, so text that
    # does not fit is simply not drawn. Extracting the rendered text back out of the PDF is
    # the only way to prove nothing was lost.
    retained = []
    if native_ran and nat.get("sb", {}).get("status") == B.NATIVE_COMPLETED:
        rendered = B.native_text(nat["sb"]["pdf"])
        retained = [r["row_id"] for r in ex["rows"]
                    if (r["controlled_display_text"] or "").strip()
                    and " ".join(r["controlled_display_text"].split()) not in rendered]
    G.append(_g("EVERY_CONTROLLED_ROW_SURVIVES_INTO_THE_RENDERED_PDF", "NATIVE",
                retained, [], len([r for r in ex["rows"]
                                   if (r["controlled_display_text"] or "").strip()]),
                "text extracted from the native PDF", NATIVE_RENDER,
                empty_by_design=not native_ran))
    G.append(_g("THE_PACKAGE_RENDER_STATUS_AGREES_WITH_EVERY_FILE_RECORD", "NATIVE",
                B.render_status(native),
                (B.NATIVE_COMPLETED
                 if native_ran and all(r["status"] == B.NATIVE_COMPLETED
                                       for r in nat.values())
                 else B.NATIVE_UNAVAILABLE if not native_ran
                 else B.NATIVE_FAILED), 1, "runtime detection", NATIVE_RENDER))
    G.append(_g("THE_RENDER_STATUS_IS_NOT_A_HARDCODED_CONSTANT", "NATIVE",
                hasattr(B, "RENDER_STATUS"), False, 1, "k5_calib_build_v1", DOC_LITERAL))

    # ------------------------------------------------------- pagination determinism -------
    G.append(_g("THE_FROZEN_METRICS_MATCH_THE_FONT_THEY_WERE_MEASURED_FROM", "DETERMINISM",
                {k: v for k, v in MT.verify_against_font().items()
                 if k in ("sha256_matches", "drifted")},
                dict(sha256_matches=True, drifted=[]), len(MT.ADVANCE),
                "k5_calib_metrics_v1 vs the font on disk", INDEPENDENT_MEASURE))
    G.append(_g("THE_FROZEN_METRICS_ARE_THE_DECLARED_VERSION_AND_FONT", "DETERMINISM",
                (MT.VERSION, MT.FONT_SHA256),
                (DECLARED_METRICS_VERSION, DECLARED_FONT_SHA256), 1,
                "k5_calib_metrics_v1", DOC_LITERAL))
    G.append(_g("THE_LINE_PITCH_IS_THE_DECLARED_FROZEN_LITERAL", "DETERMINISM",
                B.LINE_FACTOR, DECLARED_LINE_FACTOR, 1, "k5_calib_build_v1", DOC_LITERAL))
    G.append(_g("THE_PAGE_BREAK_DECISION_PATH_TOUCHES_NO_RENDERER_OR_RUNTIME_FONT",
                "DETERMINISM", sorted(_decision_path_forbidden_calls()), [], 4,
                "source of paginate/block_height_pt/wrapped_lines_pt/MT.wrap",
                INDEPENDENT_MEASURE))
    # B03-SPECIFIC golden master. Not a generator invariant: another unit will legitimately
    # produce a different map. Update it only on an approved source change, an approved
    # layout-contract change, or an explicit review decision — never to silence a failure.
    G.append(_g("B03_PAGE_BREAK_MAP_MATCHES_GOLDEN_MASTER", "GOLDEN_MASTER",
                B.page_break_map(), DECLARED_PAGE_BREAK_MAP, len(DECLARED_PAGE_BREAK_MAP),
                "review-page roll", DOC_LITERAL))

    # ------------------------------------------------- reusable relationship invariants ---
    # These hold for ANY unit. No review-page count appears in them.
    screen_ids = [s["screen_id"] for s in M.screens()]
    pages = B.slides()
    page_ids = [sl["screen"]["screen_id"] for sl in pages]
    G.append(_g("LEARNER_SCREEN_IDS_ARE_DERIVED_INDEPENDENTLY_OF_THE_REVIEW_PAGES",
                "RELATIONSHIP", screen_ids, _expected_screen_ids(pol), len(screen_ids),
                "k5_policy_apply_v1 + UNIT_ID", COMMITTED_MODEL))
    G.append(_g("EVERY_LEARNER_SCREEN_HAS_AT_LEAST_ONE_REVIEW_PAGE", "RELATIONSHIP",
                sorted(set(screen_ids) - set(page_ids)), [], len(screen_ids),
                "review-page roll", INDEPENDENT_MEASURE))
    G.append(_g("THE_REVIEW_PAGE_SCREEN_SET_EQUALS_THE_LEARNER_SCREEN_SET", "RELATIONSHIP",
                sorted(set(page_ids)), sorted(set(screen_ids)), len(set(page_ids)),
                "review-page roll", INDEPENDENT_MEASURE))
    G.append(_g("NO_REVIEW_PAGE_REFERENCES_AN_UNKNOWN_SCREEN", "RELATIONSHIP",
                sorted(set(page_ids) - set(screen_ids)), [], len(page_ids),
                "review-page roll", INDEPENDENT_MEASURE))
    seq = {}
    for sl in pages:
        seq.setdefault(sl["screen"]["screen_id"], []).append(
            (sl["page_index"], sl["page_count"]))
    incomplete = [sid for sid, v in seq.items()
                  if sorted(i for i, _ in v) != list(range(1, len(v) + 1))
                  or {c for _, c in v} != {len(v)}]
    G.append(_g("EVERY_MULTIPAGE_SEQUENCE_IS_COMPLETE_ONE_TO_N", "RELATIONSHIP",
                incomplete, [], len(seq), "review-page roll", INDEPENDENT_MEASURE))
    idents = [(sl["screen"]["screen_id"], sl["page_index"]) for sl in pages]
    G.append(_g("NO_DUPLICATE_SCREEN_ID_AND_PAGE_NUMBER_IDENTITY", "RELATIONSHIP",
                len(idents) - len(set(idents)), 0, len(idents),
                "review-page roll", INDEPENDENT_MEASURE))

    # ------------------------------------------------ canonical measured/render parity ----
    # By VALUE, not by grep. A build is replayed with the trace on; the string the
    # measurement path received and the string the writer emitted are compared to each other,
    # to the canonical text, and to what came back out of the .pptx.
    par = _parity_report()
    G.append(_g("EVERY_STORYBOARD_BLOCK_IS_MEASURED_AND_RENDERED_FROM_THE_SAME_"
                "CANONICAL_TEXT", "PARITY",
                dict(missing=par["missing"], duplicate=par["duplicate"],
                     mismatched=par["mismatched"]),
                dict(missing=[], duplicate=[], mismatched=[]), par["population"],
                f"runtime trace vs {M.STORYBOARD_NAME} readback", RAW_FILE))

    # ------------------------------------------------------------- full provenance --------
    all_ids = {r["row_id"] for r in ex["rows"]}
    notes_ids = set(re.findall(r"T03B03-ROW-\d{3}", sb["all_notes"]))
    G.append(_g("EVERY_CONTROLLED_ROW_ID_IS_MACHINE_READABLE_IN_THE_SPEAKER_NOTES",
                "TRACEABILITY", sorted(all_ids - notes_ids), [], len(all_ids),
                f"{M.STORYBOARD_NAME} notesSlide", RAW_FILE))
    return G


_FORBIDDEN_IN_DECISION_PATH = ("soffice", "libreoffice", "fitz", "ImageFont", "truetype",
                               "getlength", "textlength", "_font(", "_draw(",
                               "native_render", "native_check")


def _decision_path_forbidden_calls():
    """Read the SOURCE of every function on the page-break decision path and flag any
    renderer or runtime-font call. A claim of determinism should be checkable, not asserted."""
    import inspect
    out = []
    for fn in (B.paginate, B.block_height_pt, B.wrapped_lines_pt, MT.wrap):
        src = inspect.getsource(fn)
        body = "\n".join(ln for ln in src.splitlines()
                          if not ln.strip().startswith("#"))
        for tok in _FORBIDDEN_IN_DECISION_PATH:
            if tok in body:
                out.append(f"{fn.__module__}.{fn.__name__}:{tok}")
    return out


DECLARED_PAGE_BREAK_MAP = {
    "T03B03-S01": ["1/1"], "T03B03-S02": ["1/1"], "T03B03-S03": ["1/1"],
    "T03B03-S04": ["1/2", "2/2"], "T03B03-S05": ["1/2", "2/2"],
    "T03B03-S06": ["1/1"], "T03B03-S07": ["1/2", "2/2"], "T03B03-S08": ["1/1"],
    "T03B03-S09": ["1/2", "2/2"], "T03B03-S10": ["1/1"],
    "T03B03-S11": ["1/2", "2/2"], "T03B03-S12": ["1/1"],
}


def run(verbose=True, native=None):
    G = gates(native)
    bad = [g for g in G if not g["passed"]]
    if verbose:
        print(f"SUITE_ID = {SUITE_ID}")
        print(f"{len(G) - len(bad)}/{len(G)} passed")
        kinds = {}
        for g in G:
            kinds[g["gate_type"]] = kinds.get(g["gate_type"], 0) + 1
        print("by type:", dict(sorted(kinds.items())))
        bases = {}
        for g in G:
            bases[g["oracle_basis"]] = bases.get(g["oracle_basis"], 0) + 1
        print("by oracle:", dict(sorted(bases.items())))
        for g in bad:
            print(f"  FAIL {g['name']}\n       expected={g['expected']!r}\n"
                  f"       observed={g['observed']!r}")
    return G, bad


# ==========================================================================================
# MUTATIONS — each fixture breaks one thing; the suite must notice
# ==========================================================================================
def _rebuild(tmp, native=False):
    """Rebuild all three files into a scratch directory and point the gates at them."""
    for f in os.listdir(tmp):
        fp = os.path.join(tmp, f)
        if os.path.isfile(fp):
            os.remove(fp)
        else:
            shutil.rmtree(fp, ignore_errors=True)
    M.reset()
    old = M.PPTX_DIR
    M.PPTX_DIR = tmp
    try:
        B.build_storyboard()
        B.build_lampiran(2)
        B.build_lampiran(3)
        B.emit_handoff()
        # A native render costs ~10s per rebuild, so it runs only for the fixtures that
        # target a native gate. The rest leave those gates skipped by design.
        # (8) native output lives INSIDE this fixture's scratch, so no PDF or raster from
        # another fixture or another iteration can be read back by mistake.
        nat = (B.native_check(os.path.join(tmp, "_native"), DECLARED_BANDS_PT)
               if native else None)
        return run(verbose=False, native=nat)
    finally:
        M.PPTX_DIR = old


FIXTURES = []


def fixture(fid, target, needs_native=False):
    def deco(fn):
        FIXTURES.append(dict(id=fid, target=target, fn=fn, needs_native=needs_native))
        return fn
    return deco


def _patch(obj, key, value):
    """Set and return an undo callable, for dicts and modules alike."""
    if isinstance(obj, dict):
        had, old = key in obj, obj.get(key)
        obj[key] = value
        return lambda: (obj.__setitem__(key, old) if had else obj.pop(key, None))
    old = getattr(obj, key)
    setattr(obj, key, value)
    return lambda: setattr(obj, key, old)


# --- model-level defects -------------------------------------------------------------------
@fixture("KC-01", "EVERY_CONTROLLED_SOURCE_ROW_APPEARS_ON_THE_STORYBOARD")
def _f1():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        for s in scs:
            if s["kind"] == "CONTENT":
                s["blocks"] = [b for b in s["blocks"] if not b.get("revealed")]
                break
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KC-02", "EVERY_VERBATIM_BLOCK_EQUALS_ITS_CONTROLLED_ROW")
def _f2():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        for s in scs:
            for b in s["blocks"]:
                if b["provenance"] == "SOURCE_CONTROLLED":
                    b["text"] = b["text"] + " (nota tambahan)"
                    return scs
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KC-03", "NO_CONDENSED_LINE_ADDS_A_WORD_TO_ITS_PROPOSITION")
def _f3():
    return _patch(M, "condense", lambda t: "Kontraktor mesti mematuhi semua piawaian.")


@fixture("KC-04", "NO_BLOCK_CARRIES_THE_FORBIDDEN_PROVENANCE")
def _f4():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        scs[3]["blocks"].append(dict(text="Amalan terbaik industri.",
                                     provenance=M.FORBIDDEN_PROVENANCE))
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KC-05", "EVERY_BLOCK_USES_A_DECLARED_PROVENANCE_CLASS")
def _f5():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        scs[3]["blocks"][0] = dict(scs[3]["blocks"][0], provenance="SOMETHING_ELSE")
        return scs
    return _patch(M, "_screens_uncached", broken)


# --- state and density defects --------------------------------------------------------------
@fixture("KS-01", "RUNTIME_STATE_COUNT_MATCHES_THE_DECLARED_ROLL")
def _s1():
    orig = M._states_uncached
    return _patch(M, "_states_uncached", lambda: orig()[:-1])


@fixture("KS-02", "EVERY_PANEL_STATE_APPEARS_IN_THE_2_PANEL_LAMPIRAN")
def _s2():
    orig = M._lampiran_pages_uncached
    return _patch(M, "_lampiran_pages_uncached",
                  lambda n: orig(n)[:-1] if n == 2 else orig(n))


@fixture("KS-03", "NO_2_PANEL_PAGE_MIXES_TWO_SCREENS_OR_EXCEEDS_ITS_DENSITY")
def _s3():
    orig = M._lampiran_pages_uncached

    def broken(n):
        pages = orig(n)
        if n == 2 and len(pages) > 1:
            pages[0] = dict(pages[0],
                            states=pages[0]["states"] + pages[1]["states"][:1],
                            state_ids=pages[0]["state_ids"] + pages[1]["state_ids"][:1])
        return pages
    return _patch(M, "_lampiran_pages_uncached", broken)


@fixture("KS-04", "BOTH_DENSITIES_CARRY_THE_SAME_STATES_ONLY_THE_LAYOUT_DIFFERS")
def _s4():
    orig = M._lampiran_pages_uncached
    return _patch(M, "_lampiran_pages_uncached",
                  lambda n: orig(n)[:-2] if n == 3 else orig(n))


@fixture("KS-05", "LAMPIRAN_3_PANEL_SLIDE_COUNT_MATCHES_THE_DECLARED_ROLL")
def _s5():
    orig = M._lampiran_pages_uncached
    return _patch(M, "_lampiran_pages_uncached",
                  lambda n: orig(2) if n == 3 else orig(n))


# --- scenario defects -----------------------------------------------------------------------
@fixture("KN-01", "NO_PROPER_NAME_WAS_INVENTED_FOR_THE_CAST")
def _n1():
    return _patch(M, "PLACEHOLDER_CAST", "WATAK: Encik Rahman dan Alya")


@fixture("KN-02", "THE_SCENARIO_CARRIES_A_PENDING_CAST_PLACEHOLDER_ON_THE_FACE_OF_THE_DECK")
def _n2():
    return _patch(M, "PLACEHOLDER_CAST", "WATAK: dua pekerja tapak")


@fixture("KN-03", "EVERY_SCENARIO_LINE_IS_A_CONTROLLED_ROW_VERBATIM")
def _n3():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        s = next(x for x in scs if x["kind"] == "SCENARIO")
        s["scenario"]["items"][0]["text"] = ("Kontraktor perlu berbincang dengan jurutera "
                                             "sebelum memulakan kerja.")
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KN-04", "A_SCENARIO_SCREEN_EXISTS_AS_D3_REQUIRES")
def _n4():
    orig = M._screens_uncached

    def broken():
        return [s for s in orig() if s["kind"] != "SCENARIO"]
    return _patch(M, "_screens_uncached", broken)


# --- Rumusan defects ------------------------------------------------------------------------
@fixture("KR-01", "EVERY_MONTAGE_COMPONENT_POINTS_AT_A_COMMITTED_UNIT_BEAT")
def _r1():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        s = next(x for x in scs if x["kind"] == "RUMUSAN")
        s["rumusan"]["montage"][2]["beat"] = "Kepentingan kerja berpasukan di tapak."
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KR-02", "THE_MONTAGE_USES_A6_THREE_LABELS_IN_ORDER")
def _r2():
    return _patch(M, "A6_LABELS", ["Pengenalan", "Kandungan", "Penutup"])


@fixture("KR-03", "THE_MONTAGE_INVENTS_NO_FOURTH_SUBJECT")
def _r3():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        s = next(x for x in scs if x["kind"] == "RUMUSAN")
        s["rumusan"]["montage"].append(dict(component_id="MT-04", label="Tambahan",
                                            beat=s["rumusan"]["beats"][0],
                                            points_at_beat_index=1,
                                            marks=list(M.MONTAGE_MARKS)))
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KR-04", "THE_MONTAGE_IS_MARKED_PROVISIONAL_ON_THE_DECK")
def _r4():
    return _patch(M, "MONTAGE_MARKS", ["VISUAL_DIRECTION"])


# --- narrator, quiz, closing ------------------------------------------------------------------
@fixture("KQ-01", "THE_NARRATOR_ON_THE_DECK_IS_THE_ONE_D2_NAMES")
def _q1():
    orig = M._policy_unit

    def broken():
        u = dict(orig())
        u["narrator"] = dict(u["narrator"], name="Aiman")
        return u
    return _patch(M, "_policy_unit", broken)


@fixture("KQ-02", "EVERY_COMMITTED_STEM_APPEARS_ON_THE_DECK")
def _q2():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        s = next(x for x in scs if x["kind"] == "KUIZ")
        s["blocks"] = [b for b in s["blocks"] if not b["text"].startswith("Q3 ·")]
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KQ-03", "THE_ANSWER_KEY_IS_SHOWN_AS_DRAFTED_NOT_APPROVED")
def _q3():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        s = next(x for x in scs if x["kind"] == "KUIZ")
        for b in s["blocks"]:
            if b["text"].startswith("KUNCI JAWAPAN"):
                b["text"] = "KUNCI JAWAPAN DISAHKAN"
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KQ-04", "THE_PASS_MARK_ON_THE_DECK_IS_SIXTY_PERCENT")
def _q4():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        s = next(x for x in scs if x["kind"] == "KUIZ")
        for b in s["blocks"]:
            if b["text"].startswith("Markah lulus"):
                b["text"] = "Markah lulus 50 peratus"
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KQ-05", "EVERY_DISTRACTOR_IS_MARKED_AS_CAIR_DRAFTED")
def _q5():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        s = next(x for x in scs if x["kind"] == "KUIZ")
        for b in s["blocks"]:
            if b["provenance"] == "CAIR_DRAFTED_ASSESSMENT" and b.get("list_level"):
                b["provenance"] = "SOURCE_CONTROLLED"
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KT-01", "THE_CLOSING_TEXT_IS_D4S_INTERMEDIATE_BAHAGIAN_STRING")
def _t1():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        s = next(x for x in scs if x["kind"] == "TAMAT")
        s["blocks"][0] = dict(s["blocks"][0], text="Tamat Bahagian.")
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KT-02", "THE_DECK_DOES_NOT_CLOSE_THE_TOPIK")
def _t2():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        s = next(x for x in scs if x["kind"] == "TAMAT")
        s["blocks"][0] = dict(s["blocks"][0], text="Tamat Topik 3.")
        return scs
    return _patch(M, "_screens_uncached", broken)


# --- status and non-inheritance ---------------------------------------------------------------
@fixture("KA-01", "NO_FILE_CLAIMS_ANY_FORM_OF_APPROVAL")
def _a1():
    return _patch(M, "STATUS_MARKS", ["BARIAH_APPROVED"])


@fixture("KA-02", "EVERY_FILE_CARRIES_EVERY_DECLARED_STATUS_MARK")
def _a2():
    return _patch(M, "STATUS_MARKS", ["DRAFT_FOR_BARIAH_REVIEW"])


@fixture("KA-03", "NO_T04_IDENTIFIER_APPEARS_IN_ANY_K5_FILE")
def _a3():
    return _patch(M, "SHORT", "T04")


@fixture("KA-04", "EVERY_STORYBOARD_SCREEN_SLIDE_CARRIES_SPEAKER_NOTES")
def _a4():
    return _patch(M, "notes_for", lambda sc: "")


@fixture("KA-05", "THE_SCREEN_COUNT_COMES_FROM_BARIAHS_POLICY_NOT_FROM_T04")
def _a5():
    orig = M._screens_uncached
    return _patch(M, "_screens_uncached", lambda: orig()[:-1])


# --- layout -----------------------------------------------------------------------------------
@fixture("KL-01", "NO_STORYBOARD_PAGE_OVERFLOWS_ITS_COLUMN")
def _l1():
    return _patch(B, "paginate", lambda blocks: [list(blocks)])


@fixture("KL-02", "NO_2_PANEL_PAGE_OVERFLOWS_ITS_PANEL_BOX")
def _l2():
    orig = M._lampiran_pages_uncached

    def broken(n):
        pages = orig(n)
        if n == 2:
            merged = [s for pg in pages[:4] for s in pg["states"]]
            return [dict(page_no=1, screen_id=pages[0]["screen_id"], states=merged,
                         state_ids=[s["state_id"] for s in merged])] + pages[4:]
        return pages
    return _patch(M, "_lampiran_pages_uncached", broken)


@fixture("KL-04", "EVERY_PPTX_SLIDE_HAS_A_RENDERED_PREVIEW")
def _l4():
    orig = B.render_lampiran_previews
    return _patch(B, "render_lampiran_previews", lambda n: orig(n)[1:])


@fixture("KL-03", "THE_RENDER_STATUS_IS_NOT_A_HARDCODED_CONSTANT")
def _l3():
    """A hardcoded status constant reintroduced — the exact shape of the original defect."""
    setattr(B, "RENDER_STATUS", "RENDERED_AND_VERIFIED_IN_POWERPOINT")
    return lambda: delattr(B, "RENDER_STATUS")


# --- A4 and handoff ---------------------------------------------------------------------------
@fixture("KP-01", "THE_MODEL_DOES_NOT_CLAIM_A4_WAS_APPLIED")
def _p1():
    return _patch(M.A4_COMPARISON, "applied_in_this_deck", True)


@fixture("KP-02", "NO_SCREEN_ASSERTS_A_COMPARISON_TREATMENT")
def _p2():
    orig = M._screens_uncached

    def broken():
        scs = orig()
        for sc in scs:
            if sc["kind"] == "CONTENT":
                sc["treatment"] = "comparison"
                break
        return scs
    return _patch(M, "_screens_uncached", broken)


@fixture("KP-03", "THE_A4_EVIDENCE_ROWS_ARE_THE_ONES_THE_MODEL_CITES")
def _p3():
    return _patch(M.A4_COMPARISON, "supporting_rows", ["T03B03-ROW-001"])


@fixture("KP-04", "EVERY_A4_EVIDENCE_ROW_EXISTS_AND_CARRIES_ITS_COMPARATIVE_WORDING")
def _p4():
    ex, _ = _inputs()
    row = next(r for r in ex["rows"] if r["row_id"] == "T03B03-ROW-056")
    old = row["controlled_display_text"]
    row["controlled_display_text"] = "Swale adalah sejenis parit yang cetek dan lebar."

    def undo():
        row["controlled_display_text"] = old
    return undo


@fixture("KH-01", "THE_HANDOFF_LISTS_EVERY_REQUIRED_ITEM")
def _h1():
    return _patch(B, "HANDOFF_NAME", "SOMEWHERE_ELSE_v0_1.md")


@fixture("KH-02", "THE_HANDOFF_CLAIMS_NO_NATIVE_POWERPOINT_VALIDATION")
def _h2():
    return _patch(B, "RENDER_NOTE",
                  "The deck was rendered in PowerPoint and visually verified.")


# --- render divergence ------------------------------------------------------------------------
@fixture("KX-01", "EVERY_BOUNDED_TEXT_BOX_EXPLICITLY_DECLARES_NO_AUTOFIT")
def _x1():
    """spAutoFit reintroduced: the renderer may grow the box past its declared height."""
    return _patch(B, "strip_autofit", lambda prs: 0)


@fixture("KX-08", "EVERY_BOUNDED_TEXT_BOX_EXPLICITLY_DECLARES_NO_AUTOFIT")
def _x8():
    """The autofit child is REMOVED but noAutofit is never written.

    The old negative gate passed this: nothing forbidden was present. The package still
    declared no contract, which is the ambiguity the positive gate exists to reject.
    """
    from pptx.oxml.ns import qn

    def strip_only(tf):
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        for tag in ("a:spAutoFit", "a:normAutofit", "a:noAutofit"):
            for el in bodyPr.findall(qn(tag)):
                bodyPr.remove(el)
        return tf
    undo_a = _patch(B, "_no_autofit", strip_only)

    def strip_all(prs):
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    strip_only(shape.text_frame)
        return 0
    undo_b = _patch(B, "strip_autofit", strip_all)

    def undo():
        undo_a()
        undo_b()
    return undo


@fixture("KX-02", "NO_RENDERED_INK_ENTERS_THE_FOOTER_BAND", needs_native=True)
def _x2():
    """Review pagination omitted: every screen becomes one page and long screens spill."""
    return _patch(B, "paginate", lambda blocks: [list(blocks)])


@fixture("KX-03", "NO_RENDERED_INK_ENTERS_THE_FOOTER_BAND", needs_native=True)
def _x3():
    """Body limit pushed past the footer: rendered text enters the footer band."""
    return _patch(B, "BODY_LIMIT_PT", 520.0)


@fixture("KX-04", "THE_NATIVE_RENDER_COMPLETED_FOR_EVERY_FILE", needs_native=True)
def _x4():
    """Renderer falsely reported unavailable: the gate must not accept a silent skip."""
    return _patch(B, "libreoffice", lambda: (None, None))


@fixture("KX-05", "THE_PACKAGE_RENDER_STATUS_AGREES_WITH_EVERY_FILE_RECORD",
         needs_native=True)
def _x5():
    """A failed native render reported as a pass."""
    orig = B.native_render

    def broken(pptx_path, outdir):
        rec = orig(pptx_path, outdir)
        if "STORYBOARD" in pptx_path:
            rec = dict(rec, status=B.NATIVE_FAILED)
        return rec
    undo_render = _patch(B, "native_render", broken)
    undo_status = _patch(B, "render_status", lambda native=None: B.NATIVE_COMPLETED)

    def undo():
        undo_render()
        undo_status()
    return undo


@fixture("KX-06", "NO_RENDERED_INK_LEAVES_THE_PAGE_CANVAS", needs_native=True)
def _x6():
    """Body limit pushed past the page itself: rendered ink exits the canvas at the bottom
    edge. Distinct from footer entry — the page is 540pt and the limit becomes 900pt."""
    return _patch(B, "BODY_LIMIT_PT", 900.0)


@fixture("KX-07", "EVERY_NATIVE_PDF_PAGE_COUNT_EQUALS_ITS_PPTX_SLIDE_COUNT",
         needs_native=True)
def _x7():
    """NATIVE PDF PAGE-LOSS mutation.

    The renderer silently drops a slide and the PDF comes back one page short. This is a
    page-loss fixture, NOT a stale-output fixture. Stale-output protection is a property of
    the harness, not of any single fixture: unique random scratch per fixture, empty-directory
    assertion before execution, a fixture-scoped `_native` directory, a clean rebuild that
    must pass before the mutation is applied, no PDF or raster reuse across fixtures, cleanup
    on success, and retain-on-failure with the path reported.
    """
    orig = B.native_render

    def broken(pptx_path, outdir):
        rec = orig(pptx_path, outdir)
        if rec["status"] == B.NATIVE_COMPLETED and "STORYBOARD" in pptx_path:
            import fitz
            with fitz.open(rec["pdf"]) as d:
                d.delete_page(1)
                d.save(rec["pdf"] + ".short")
            os.replace(rec["pdf"] + ".short", rec["pdf"])
            with fitz.open(rec["pdf"]) as d:
                rec["page_count"] = d.page_count
        return rec
    return _patch(B, "native_render", broken)


# --- determinism and retention ----------------------------------------------------------------
@fixture("KD-01", "THE_LINE_PITCH_IS_THE_DECLARED_FROZEN_LITERAL")
def _d1():
    return _patch(B, "LINE_FACTOR", 1.0)


@fixture("KD-02", "THE_PAGE_BREAK_DECISION_PATH_TOUCHES_NO_RENDERER_OR_RUNTIME_FONT")
def _d2():
    """Runtime font measurement reintroduced into the decision path."""
    orig = MT.wrap

    def measured(text, size_pt, width_pt):
        from PIL import ImageFont          # noqa: F401
        return orig(text, size_pt, width_pt)
    return _patch(MT, "wrap", measured)


@fixture("KD-03", "B03_PAGE_BREAK_MAP_MATCHES_GOLDEN_MASTER")
def _d3():
    return _patch(B, "BODY_LIMIT_PT", 300.0)


@fixture("KV-01", "EVERY_CONTROLLED_ROW_SURVIVES_INTO_THE_RENDERED_PDF", needs_native=True)
def _v1():
    """A block silently dropped after pagination — the clipping this gate exists to catch."""
    orig = B.paginate

    def broken(blocks):
        pages = orig(blocks)
        if len(pages[0]) > 3:
            pages[0] = pages[0][:-2]
        return pages
    return _patch(B, "paginate", broken)


@fixture("KV-02", "EVERY_CONTROLLED_ROW_ID_IS_MACHINE_READABLE_IN_THE_SPEAKER_NOTES")
def _v2():
    """Notes truncated: the complete row-ID set stops being machine-verifiable."""
    orig = M.notes_for
    return _patch(M, "notes_for",
                  lambda sc: orig(sc).split("Baris sumber:")[0])


RETAINED_ON_FAILURE = []


@fixture("KY-01", "EVERY_STORYBOARD_BLOCK_IS_MEASURED_AND_RENDERED_FROM_THE_SAME_"
                  "CANONICAL_TEXT")
def _y1():
    """Recreates the discovered defect on ONE path only: the writer re-inserts the ' · '
    separator that pagination does not measure."""
    orig = B.block_text
    return _patch(B, "render_text",
                  lambda b: orig(b).replace("] ", "] · ", 1))


# ---- registry integrity ----------------------------------------------------------------------
def registry_integrity(executed_ids=None):
    """Fixture accounting derived FROM the registry, never hand-written."""
    import collections
    ids = [f["id"] for f in FIXTURES]
    dup = [k for k, v in collections.Counter(ids).items() if v > 1]
    ex = list(executed_ids or [])
    return dict(registry_count=len(ids), ordered_ids=ids, unique=not dup, duplicates=dup,
                native_fixtures=sum(1 for f in FIXTURES if f.get("needs_native")),
                executed_count=len(ex),
                registered_not_executed=sorted(set(ids) - set(ex)),
                executed_not_registered=sorted(set(ex) - set(ids)),
                reconciles=(not dup and sorted(ids) == sorted(ex)) if ex else None)


def run_mutations(verbose=True):
    base_nat = B.native_check(NATIVE_DIR, DECLARED_BANDS_PT)
    base_G, base_bad = run(verbose=False, native=base_nat)
    results, missed, dirty = [], [], []
    root = tempfile.mkdtemp(prefix="k5calib_root_")
    # The layout gates need the MEASUREMENT, not the PNGs; writing 48 images per fixture is
    # pure I/O. Measurement is unchanged, so KL-01 and KL-02 still bite.
    B.WRITE_PREVIEW_IMAGES = False
    try:
        for fx in FIXTURES:
            # (1) unique fresh scratch per fixture — never shared, never reused
            tmp = tempfile.mkdtemp(prefix=f"k5fx_{fx['id']}_", dir=root)
            nat_dir = os.path.join(tmp, "_native")
            # (2) empty-directory assertion, before anything is written
            assert os.listdir(tmp) == [], f"{fx['id']}: scratch not empty at start"
            keep = False
            try:
                # (3)+(4) clean direction: the UNMUTATED build must pass in this same
                # directory first. A stale PDF, raster or manifest leaking in from another
                # fixture shows up here as a clean-run failure, not as a false detection.
                _, clean_bad = _rebuild(tmp, native=fx.get("needs_native", False))
                clean_ok = not clean_bad
                if not clean_ok:
                    dirty.append(dict(id=fx["id"],
                                      failed=sorted(g["name"] for g in clean_bad)))
                # (5) exactly one declared mutation
                undo = fx["fn"]()
                try:
                    # (6) rerun; (7) assert the SPECIFIC expected gate fails
                    G, bad = _rebuild(tmp, native=fx.get("needs_native", False))
                    names = {g["name"] for g in bad}
                    caught = fx["target"] in names
                    results.append(dict(id=fx["id"], target=fx["target"], detected=caught,
                                        clean_before=clean_ok,
                                        scratch=tmp, native_dir=nat_dir,
                                        also_failed=sorted(names - {fx["target"]})))
                    if not caught or not clean_ok:
                        missed.append(fx["id"])
                        keep = True
                finally:
                    undo()
                    M.reset()
            finally:
                # (9) cleaned, or retained on failure with the path reported
                if keep:
                    RETAINED_ON_FAILURE.append(tmp)
                else:
                    shutil.rmtree(tmp, ignore_errors=True)
    finally:
        B.WRITE_PREVIEW_IMAGES = True
        if not RETAINED_ON_FAILURE:
            shutil.rmtree(root, ignore_errors=True)
    # Rebuild the real files so the on-disk artifacts match the unmutated model.
    M.reset()
    B.build_storyboard()
    B.build_lampiran(2)
    B.build_lampiran(3)
    reg = registry_integrity([r["id"] for r in results])
    out = dict(baseline_pass=len(base_G) - len(base_bad), baseline_total=len(base_G),
               baseline_executions=1,
               fixture_count=reg["registry_count"],
               detected=reg["registry_count"] - len(missed),
               missed=missed,
               clean_before_executions=len(results),
               mutated_executions=len(results),
               total_suite_executions=1 + 2 * len(results),
               native_fixtures=reg["native_fixtures"],
               registry=reg,
               clean_before_all_passed=not dirty, dirty=dirty,
               retained_on_failure=RETAINED_ON_FAILURE, results=results)
    assert reg["reconciles"], f"fixture registry does not reconcile: {reg}"
    if verbose:
        print(json.dumps({k: v for k, v in out.items() if k != "results"}, indent=1))
        for r in results:
            print(f"  {r['id']}  {'OK ' if r['detected'] else 'MISS'}  {r['target']}")
    return out


if __name__ == "__main__":
    if "--mutations" in sys.argv:
        r = run_mutations()
        sys.exit(0 if not r["missed"] else 1)
    nat = B.native_check(NATIVE_DIR, DECLARED_BANDS_PT)
    G, bad = run(native=nat)
    print(f"native: {B.render_status(nat)}")
    sys.exit(0 if not bad else 1)
